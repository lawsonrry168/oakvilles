# -*- coding: utf-8 -*-
"""Generate client-facing before/after comparison deck (simplified PPTX)."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import datetime
import os
import shutil

ROOT = os.path.join(r"c:\Users\Steriod\Desktop\oakvilles")
OUTPUT = os.path.join(ROOT, "Oakville-Site-Comparison-Client.pptx")
OUTPUT_PRESENTATIONS = os.path.join(ROOT, "presentations", "Oakville-Site-Comparison-Client.pptx")

PINE = RGBColor(0x2A, 0x46, 0x3C)
PINE_LIGHT = RGBColor(0x3D, 0x5A, 0x4E)
PAPER = RGBColor(0xF7, 0xF2, 0xE7)
PAPER_DARK = RGBColor(0xE8, 0xE0, 0xD0)
CINNABAR = RGBColor(0xA2, 0x3A, 0x2E)
OCHRE = RGBColor(0xAE, 0x8A, 0x47)
INK = RGBColor(0x1A, 0x1A, 0x18)
MUTED = RGBColor(0x5A, 0x5A, 0x52)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OLD_GRAY = RGBColor(0x9A, 0x90, 0x88)

FONT = "Microsoft JhengHei"
today = datetime.date(2026, 6, 29)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_fill(shape, color: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color: RGBColor):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    set_fill(s, color)
    s.line.fill.background()
    return s


def add_textbox(slide, left, top, width, height, text, size=14, bold=False,
                color=INK, align=PP_ALIGN.LEFT, line_spacing=1.2):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.line_spacing = line_spacing
    r = p.runs[0]
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def add_bullets(slide, left, top, width, height, items, size=12, color=INK, spacing=Pt(5)):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = spacing
        p.line_spacing = 1.12
        r = p.runs[0]
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return tb


def slide_header(slide, section_num: str, title: str, subtitle: str = ""):
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, PAPER)
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), CINNABAR)
    add_rect(slide, 0, Inches(0.08), SLIDE_W, Inches(1.05), PINE)
    if section_num:
        add_textbox(slide, Inches(0.55), Inches(0.22), Inches(1.2), Inches(0.35),
                    section_num, 11, True, OCHRE)
    add_textbox(slide, Inches(0.55), Inches(0.48), Inches(12), Inches(0.55),
                title, 26, True, WHITE)
    if subtitle:
        add_textbox(slide, Inches(0.55), Inches(0.95), Inches(12), Inches(0.35),
                    subtitle, 11, False, PAPER_DARK)
    add_textbox(slide, Inches(0.55), Inches(7.05), Inches(6), Inches(0.3),
                "頤安本草 · 官網重建前後比較", 9, False, MUTED)
    add_textbox(slide, Inches(10.5), Inches(7.05), Inches(2.3), Inches(0.3),
                today.strftime("%Y-%m-%d"), 9, False, MUTED, PP_ALIGN.RIGHT)


def stat_card(slide, left, top, number, label, w=Inches(2.55), h=Inches(1.45)):
    add_rect(slide, left, top, w, h, WHITE)
    add_rect(slide, left, top, Inches(0.06), h, CINNABAR)
    add_textbox(slide, left + Inches(0.18), top + Inches(0.12), w - Inches(0.28), Inches(0.65),
                number, 32, True, PINE)
    add_textbox(slide, left + Inches(0.18), top + Inches(0.78), w - Inches(0.28), Inches(0.5),
                label, 10, False, MUTED)


def comparison_col(slide, left, top, width, header, header_color, items, bg=WHITE):
    h = Inches(0.42)
    add_rect(slide, left, top, width, h, header_color)
    add_textbox(slide, left + Inches(0.12), top + Inches(0.06), width - Inches(0.24), h,
                header, 12, True, WHITE)
    body_h = Inches(4.35)
    add_rect(slide, left, top + h, width, body_h, bg)
    add_bullets(slide, left + Inches(0.12), top + h + Inches(0.12),
                width - Inches(0.24), body_h - Inches(0.15), items, 11, INK)


def before_after_row(slide, y, label, old_text, new_text):
    add_rect(slide, Inches(0.55), y, Inches(1.45), Inches(0.72), PINE)
    add_textbox(slide, Inches(0.62), y + Inches(0.14), Inches(1.3), Inches(0.45),
                label, 10, True, WHITE)
    add_rect(slide, Inches(2.05), y, Inches(5.25), Inches(0.72), RGBColor(0xF0, 0xED, 0xEA))
    add_textbox(slide, Inches(2.15), y + Inches(0.1), Inches(5.05), Inches(0.55),
                old_text, 10, False, MUTED)
    add_rect(slide, Inches(7.35), y, Inches(5.4), Inches(0.72), WHITE)
    add_rect(slide, Inches(7.35), y, Inches(0.05), Inches(0.72), CINNABAR)
    add_textbox(slide, Inches(7.5), y + Inches(0.1), Inches(5.15), Inches(0.55),
                new_text, 10, True, PINE)


prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]

# ── 1 Cover ──
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, PINE)
add_rect(s, 0, 0, Inches(0.12), SLIDE_H, CINNABAR)
add_textbox(s, Inches(0.9), Inches(1.75), Inches(8), Inches(0.45), "OAKVILLE WELLNESS", 13, False, OCHRE)
add_textbox(s, Inches(0.9), Inches(2.25), Inches(10.5), Inches(1.0),
            "頤安本草官網 · 前後比較", 38, True, WHITE)
add_textbox(s, Inches(0.9), Inches(3.35), Inches(10), Inches(0.55),
            "舊站 → 新站重建 → 最新上線版本", 20, False, PAPER)
add_textbox(s, Inches(0.9), Inches(4.15), Inches(10), Inches(0.4),
            f"客戶交付精簡版 ｜ {today.strftime('%Y 年 %m 月')} ｜ WhatsApp 6734 9532", 12, False, PAPER_DARK)

# ── 2 Summary ──
s = prs.slides.add_slide(blank)
slide_header(s, "01", "一句話結論")
add_textbox(s, Inches(0.55), Inches(1.32), Inches(12.2), Inches(0.9),
            "舊站為單頁結構，無法承接「症狀搜尋」與轉介流量。"
            "新站以 62 頁雙語架構補足 SEO、信任與 WhatsApp 轉化，"
            "並保持中環高端、合規克制的品牌定位。",
            14, False, INK)
stat_card(s, Inches(0.55), Inches(2.45), "1→62", "頁面（含雙語）", Inches(2.85))
stat_card(s, Inches(3.55), Inches(2.45), "8", "症狀 SEO 頁", Inches(2.85))
stat_card(s, Inches(6.55), Inches(2.45), "5.0", "Google 評價", Inches(2.85))
stat_card(s, Inches(9.55), Inches(2.45), "✓", "Vercel 已上線", Inches(2.85))
add_rect(s, Inches(0.55), Inches(4.2), Inches(12.2), Inches(2.5), WHITE)
add_textbox(s, Inches(0.75), Inches(4.35), Inches(11.8), Inches(0.32), "本簡報涵蓋", 13, True, PINE)
add_bullets(s, Inches(0.75), Inches(4.72), Inches(11.5), Inches(1.85), [
    "舊站與新站核心差距（一表看清）",
    "重建後 8 大改善與最新一輪潤飾重點",
    "現已上線內容、網址與建議下一步",
], 12, INK)

# ── 3 Old vs New table ──
s = prs.slides.add_slide(blank)
slide_header(s, "02", "舊站 vs 新站", "重建前 oakvilles.com ｜ 現版 oakvilles.vercel.app")
add_textbox(s, Inches(2.05), Inches(1.38), Inches(5.25), Inches(0.3), "舊站（重建前）", 11, True, OLD_GRAY)
add_textbox(s, Inches(7.35), Inches(1.38), Inches(5.4), Inches(0.3), "新站（現版）", 11, True, PINE)
y = Inches(1.72)
rows = [
    ("頁面", "實質 1 頁 · 導航 5 項", "62 頁中+英 · 30+ 核心內容頁"),
    ("品牌", "伍厚臻個人 · 訊息分散", "頤安本草 · 印章識別 · 禪意視覺"),
    ("SEO", "無 sitemap／Schema／症狀頁", "8 症狀頁 + sitemap + JSON-LD + OG"),
    ("轉化", "僅 WA 浮窗", "2 步預約 funnel · 計算器 · sticky CTA"),
    ("信任", "缺評價、環境圖、FAQ", "Google 5.0 · 診所實景 · 明碼收費"),
    ("量度", "無 GA4／CTA 追蹤", "6 個轉化事件就緒（待填 GA4）"),
    ("合規", "部分療效承諾", "合規克制 · 評價免責"),
]
for label, old_t, new_t in rows:
    before_after_row(s, y, label, old_t, new_t)
    y += Inches(0.78)

# ── 4 Eight improvements ──
s = prs.slides.add_slide(blank)
slide_header(s, "03", "新站 8 大改善")
items = [
    ("資訊架構", "1 頁 → 62 頁雙語，覆蓋專科、症狀、Blog、FAQ"),
    ("本地 SEO", "新增中環專頁 central-hk · llms.txt · hreflang"),
    ("預約動線", "2 步 funnel 自動組 WhatsApp 訊息"),
    ("收費透明", "流程收費頁 + 診金即時計算器"),
    ("社交分享", "每頁 OG 1200×630 · 品牌首頁分享圖"),
    ("搜尋技術", "sitemap · canonical · 結構化資料"),
    ("手機體驗", "sticky CTA · WA 浮動避障 · Mobile-first"),
    ("部署維護", "GitHub + Vercel 一鍵建置上線"),
]
for i, (title, desc) in enumerate(items):
    col = i % 4
    row = i // 4
    x = Inches(0.55) + col * Inches(3.15)
    y = Inches(1.48) + row * Inches(2.75)
    add_rect(s, x, y, Inches(2.95), Inches(2.45), WHITE)
    add_rect(s, x, y, Inches(2.95), Inches(0.38), PINE if row == 0 else PINE_LIGHT)
    add_textbox(s, x + Inches(0.12), y + Inches(0.06), Inches(2.7), Inches(0.28), title, 12, True, WHITE)
    add_textbox(s, x + Inches(0.12), y + Inches(0.48), Inches(2.7), Inches(1.85), desc, 10.5, False, INK)

# ── 5 Latest polish ──
s = prs.slides.add_slide(blank)
slide_header(s, "04", "最新一輪潤飾", "2026 年 6 月 · 已同步 GitHub 與 Vercel")
comparison_col(s, Inches(0.55), Inches(1.42), Inches(5.9), "早前版本", OLD_GRAY, [
    "WhatsApp 連結格式不統一",
    "部分 meta 描述過短",
    "OG 圖僅首頁為主",
    "英文 Hero 浮水印位置不穩",
    "印章為純文字四格",
    "sitemap priority 固定",
])
comparison_col(s, Inches(6.85), Inches(1.42), Inches(5.9), "現版（已完成）", PINE, [
    "全站統一 wa.me + 預約問候語模板",
    "全站 description 110–165 字",
    "default／blog／process 專用 OG 圖",
    "Heal 浮水印佈局修正 + 品牌字體",
    "PNG 印章 + 5 枚主題圖章",
    "lastmod 日期 + 分頁 priority 邏輯",
])
add_rect(s, Inches(0.55), Inches(5.95), Inches(12.2), Inches(0.85), PINE_LIGHT)
add_textbox(s, Inches(0.75), Inches(6.08), Inches(11.8), Inches(0.6),
            "新增著陸頁：中環暗瘡濕疹中醫（中英）｜ 新聞頁 Article schema ｜ 首頁無障礙 main  landmark",
            10.5, False, PAPER)

# ── 6 What's live ──
s = prs.slides.add_slide(blank)
slide_header(s, "05", "現已上線內容")
modules = [
    ("雙語官網", "中文主站 + 完整 /en/ 英文鏡像"),
    ("症狀 SEO", "濕疹·暗瘡·失眠·備孕·頸痛·坐骨 + 中環綜合頁"),
    ("專科療法", "痛症·皮膚·婦科·內科 + 針灸·中藥·艾灸·拔罐"),
    ("內容信任", "Blog 4 篇 · FAQ · 醫師簡介 · Google 5.0"),
    ("轉化工具", "WhatsApp 6734 9532 · 2 步預約 · 診金計算器"),
    ("技術基礎", "sitemap · Schema · OG · dataLayer 事件"),
]
for i, (title, desc) in enumerate(modules):
    col = i % 3
    row = i // 3
    x = Inches(0.55) + col * Inches(4.15)
    y = Inches(1.48) + row * Inches(2.35)
    add_rect(s, x, y, Inches(3.85), Inches(2.05), WHITE)
    add_rect(s, x, y, Inches(3.85), Inches(0.4), PINE)
    add_textbox(s, x + Inches(0.12), y + Inches(0.05), Inches(3.5), Inches(0.32), title, 13, True, WHITE)
    add_textbox(s, x + Inches(0.12), y + Inches(0.52), Inches(3.5), Inches(1.4), desc, 11, False, INK)

# ── 7 Deployment ──
s = prs.slides.add_slide(blank)
slide_header(s, "06", "部署與網址")
add_rect(s, Inches(0.55), Inches(1.42), Inches(5.9), Inches(2.2), WHITE)
add_textbox(s, Inches(0.75), Inches(1.55), Inches(5.5), Inches(0.32), "✓  新站（最新版）", 14, True, PINE)
add_bullets(s, Inches(0.75), Inches(1.95), Inches(5.5), Inches(1.5), [
    "網址：oakvilles.vercel.app",
    "GitHub main 已同步（ea6db4e）",
    "Vercel Production 建置成功",
    "64 頁靜態輸出 · 自動 build:deploy",
], 11.5, INK)
add_rect(s, Inches(6.85), Inches(1.42), Inches(5.9), Inches(2.2), WHITE)
add_rect(s, Inches(6.85), Inches(1.42), Inches(0.06), Inches(2.2), CINNABAR)
add_textbox(s, Inches(7.05), Inches(1.55), Inches(5.5), Inches(0.32), "⚠  自訂網域（待接）", 14, True, CINNABAR)
add_bullets(s, Inches(7.05), Inches(1.95), Inches(5.5), Inches(1.5), [
    "oakvilles.com 目前仍指向舊 hosting",
    "需將 DNS 改指向 Vercel",
    "完成後訪客自動看到新站",
    "建議本週內完成切換",
], 11.5, INK)
add_rect(s, Inches(0.55), Inches(3.85), Inches(12.2), Inches(2.65), PINE)
add_textbox(s, Inches(0.75), Inches(4.0), Inches(11.8), Inches(0.32), "共同轉化 KPI", 13, True, OCHRE)
add_textbox(s, Inches(0.75), Inches(4.45), Inches(11.5), Inches(1.85),
            "所有預約入口統一導向 WhatsApp +852 6734 9532，"
            "並附標準問候語「您好，我想預約伍醫師診症」。"
            "網站已預留 whatsapp_click 等 6 個追蹤事件，"
            "填入 GA4 後即可量度各入口成效。",
            12, False, PAPER)

# ── 8 Next steps ──
s = prs.slides.add_slide(blank)
slide_header(s, "07", "建議下一步")
cols = [
    ("P0 · 本週", PINE, [
        "oakvilles.com DNS 指向 Vercel",
        "填入 GA4 並重新部署",
        "Search Console 提交 sitemap",
        "Google Business Profile 優化",
    ]),
    ("P1 · 廣告前", OCHRE, [
        "GTM：GA4 + Meta Pixel",
        "GA4 轉化匯入 Google Ads",
        "Meta 自訂轉換 whatsapp_click",
    ]),
    ("P2 · 持續", CINNABAR, [
        "Blog 每月 1–2 篇",
        "品牌視覺資產全站接入",
        "症狀廣告著陸頁 A/B 測試",
    ]),
]
for i, (title, color, items) in enumerate(cols):
    x = Inches(0.55) + i * Inches(4.15)
    add_rect(s, x, Inches(1.45), Inches(3.85), Inches(5.1), WHITE)
    add_rect(s, x, Inches(1.45), Inches(3.85), Inches(0.42), color)
    add_textbox(s, x + Inches(0.12), Inches(1.52), Inches(3.6), Inches(0.32), title, 13, True, WHITE)
    add_bullets(s, x + Inches(0.12), Inches(2.05), Inches(3.6), Inches(4.2), items, 11.5, INK)

# ── 9 Closing ──
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, PINE)
add_rect(s, 0, 0, Inches(0.12), SLIDE_H, CINNABAR)
add_textbox(s, Inches(0.9), Inches(1.55), Inches(11), Inches(0.4), "總結", 13, False, OCHRE)
add_textbox(s, Inches(0.9), Inches(2.1), Inches(11.5), Inches(2.4),
            "舊站無法承接搜尋與轉介流量；"
            "新站已在架構、SEO、信任元件與 WhatsApp 轉化上系統補足，"
            "並完成最新一輪全站潤飾與上線部署。\n\n"
            "下一步重點：自訂網域切換 + GA4 啟用 + 本地 SEO。",
            20, False, WHITE)
add_rect(s, Inches(0.9), Inches(4.95), Inches(11.5), Inches(1.55), PINE_LIGHT)
add_textbox(s, Inches(1.1), Inches(5.12), Inches(11), Inches(0.32), "聯絡", 12, True, OCHRE)
add_textbox(s, Inches(1.1), Inches(5.5), Inches(11), Inches(0.85),
            "頤安本草 · 伍厚臻中醫師 ｜ 中環皇后大道中 99 號錦安大廈 6 樓 602 室\n"
            "WhatsApp +852 6734 9532 ｜ 電話 2881 8182 ｜ oakvilles.vercel.app",
            12, False, PAPER)

prs.save(OUTPUT)
os.makedirs(os.path.dirname(OUTPUT_PRESENTATIONS), exist_ok=True)
shutil.copy2(OUTPUT, OUTPUT_PRESENTATIONS)
print(f"Saved: {OUTPUT}")
print(f"Copied: {OUTPUT_PRESENTATIONS}")
print(f"Slides: {len(prs.slides)}")
