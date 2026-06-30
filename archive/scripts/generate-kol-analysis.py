# -*- coding: utf-8 -*-
"""Generate KOL analysis DOCX + PPTX for 養康對標 / 頤安行銷預算."""

from docx import Document
from docx.shared import Pt, Cm, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
from pptx import Presentation
from pptx.util import Inches, Pt as PptPt
from pptx.dml.color import RGBColor as PptRGB
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import datetime
import os
import shutil

BASE = r"c:\Users\Steriod\Desktop\oakvilles"
DOCX_OUT = os.path.join(BASE, "Oakville-KOL-Marketing-Analysis.docx")
DOCX_DESKTOP = r"c:\Users\Steriod\Desktop\Oakville-KOL-Marketing-Analysis.docx"
DOCX_PRES = os.path.join(BASE, "presentations", "Oakville-KOL-Marketing-Analysis.docx")
PPTX_OUT = os.path.join(BASE, "Oakville-KOL-Marketing-Analysis.pptx")
PPTX_DESKTOP = r"c:\Users\Steriod\Desktop\Oakville-KOL-Marketing-Analysis.pptx"
PPTX_PRES = os.path.join(BASE, "presentations", "Oakville-KOL-Marketing-Analysis.pptx")

today = datetime.date(2026, 6, 22)

# KOL 資料（養康 @hons_cmc 合作案例，含公開 IG 資料）
KOLS = [
    {
        "handle": "amyng_amy",
        "name": "Amy Ng",
        "followers": "1.8 萬",
        "tier": "小型 KOL",
        "niche": "美妝／生活／時尚／飲食",
        "collab": "肩頸痛症體驗文（多圖輪播）",
        "treatment": "針灸、農本方中藥",
        "likes": "9",
        "year": "2015",
        "url": "https://www.instagram.com/amyng_amy",
        "notes": "U Blog 星級博客；活躍於 IG、Facebook、小紅書等",
    },
    {
        "handle": "pinksheeb",
        "name": "Sandra Ho",
        "followers": "4,237",
        "tier": "微型 KOL",
        "niche": "美食／美容／生活",
        "collab": "中環養康到店體驗",
        "treatment": "拔火罐、五行艾灸（韓製艾灸器）",
        "likes": "9",
        "year": "2015",
        "url": "https://www.instagram.com/pinksheeb",
    },
    {
        "handle": "meikiw",
        "name": "黃美棋 Meiki Wong",
        "followers": "11 萬（認證）",
        "tier": "中階＋藝人",
        "niche": "藝人／主持／美容健康",
        "collab": "養康中醫館針灸體驗",
        "treatment": "腰痛、腎虛調理、背部針灸",
        "likes": "1,091",
        "year": "2010s",
        "url": "https://www.instagram.com/meikiw",
    },
    {
        "handle": "eva_pinkland",
        "name": "Eva Cheung",
        "followers": "26.2 萬（認證）",
        "tier": "大型 KOL",
        "niche": "時尚／創意／UK×HK",
        "collab": "標記 @hons_cmc",
        "treatment": "中環養康首次針灸、腸胃調理＋5 日中藥",
        "likes": "5,806",
        "year": "2010s",
        "url": "https://www.instagram.com/eva_pinkland",
    },
    {
        "handle": "stephanie_ssw",
        "name": "申倩瑋 Stephanie Sun",
        "followers": "5.5 萬（認證）",
        "tier": "中階 KOL",
        "niche": "美容／旅遊／生活",
        "collab": "標記 @hons_cmc",
        "treatment": "首次針灸、肚/手/腳、調理體質",
        "likes": "3,320",
        "year": "2010s",
        "url": "https://www.instagram.com/stephanie_ssw",
    },
    {
        "handle": "hanatam",
        "name": "譚杏藍 Hana Tam",
        "followers": "40.6 萬（認證）",
        "tier": "大型 KOL",
        "niche": "生活／飲食／時尚",
        "collab": "標記 @hons_cmc",
        "treatment": "印堂針灸、改善睡眠",
        "likes": "7,243",
        "year": "2010s",
        "url": "https://www.instagram.com/hanatam",
    },
    {
        "handle": "cherrietam",
        "name": "Cherrie Tam",
        "followers": "5.8 萬（認證）",
        "tier": "中階 KOL",
        "niche": "美妝／育兒／創業",
        "collab": "多次針灸分享、提及養康",
        "treatment": "針灸助眠、養康肚部針灸（腸胃）",
        "likes": "430–678",
        "year": "2010s",
        "url": "https://www.instagram.com/cherrietam",
    },
    {
        "handle": "bithialam",
        "name": "Bithia Lam",
        "followers": "4.3 萬（認證）",
        "tier": "中階 KOL",
        "niche": "媽媽／家庭／生活",
        "collab": "提及養康中醫師",
        "treatment": "針灸＋中藥調理濕疹（3 日改善分享）",
        "likes": "366",
        "year": "2010s",
        "url": "https://www.instagram.com/bithialam",
    },
]

PRICE_TIERS = [
    ("微型", "< 1 萬粉絲", "800 – 2,500", "1,500 – 5,000", "4,000 – 12,000"),
    ("小型", "1 – 5 萬", "2,500 – 8,000", "5,000 – 15,000", "12,000 – 36,000"),
    ("中階", "5 – 15 萬", "8,000 – 25,000", "15,000 – 40,000", "36,000 – 100,000"),
    ("大型", "15 – 50 萬", "25,000 – 80,000", "40,000 – 120,000", "100,000 – 300,000"),
    ("頂級", "50 萬以上", "80,000 – 300,000+", "120,000 – 500,000+", "300,000 – 1,000,000+"),
]

BUDGET_PLANS = [
    ("A · 輕量", "4 次/年", "小型 KOL × 4（單圖帖）", "約 HKD 24,000", "月均 ~2,000"),
    ("B · 標準", "6 次/年", "中階 × 2 + 小型 × 4", "約 HKD 64,000", "月均 ~5,300"),
    ("C · 加強", "8 次/年", "大型 × 1 + 中階 × 2 + 小型 × 3", "約 HKD 130,000", "月均 ~10,800"),
    ("D · 對標養康", "4–6 次/年", "大型 × 2 + 中階 × 2（如 eva + hanatam 級）", "約 HKD 180,000+", "預算較高，按需選用"),
]

PINE = PptRGB(0x2A, 0x46, 0x3C)
PINE_LIGHT = PptRGB(0x3D, 0x5A, 0x4E)
PAPER = PptRGB(0xF7, 0xF2, 0xE7)
CINNABAR = PptRGB(0xA2, 0x3A, 0x2E)
OCHRE = PptRGB(0xAE, 0x8A, 0x47)
INK = PptRGB(0x1A, 0x1A, 0x18)
MUTED = PptRGB(0x5A, 0x5A, 0x52)
WHITE = PptRGB(0xFF, 0xFF, 0xFF)
BENCH = PptRGB(0x6B, 0x8F, 0x7A)
FONT = "Microsoft JhengHei"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ── DOCX helpers ──

def set_cell_shading(cell, fill: str):
    shading = OxmlElement("w:shd")
    shading.set(qn("w:fill"), fill)
    shading.set(qn("w:val"), "clear")
    cell._tc.get_or_add_tcPr().append(shading)


def style_run(run, size=10.5, bold=False, color=None):
    run.bold = bold
    run.font.name = "Arial"
    run._element.rPr.rFonts.set(qn("w:eastAsia"), "Microsoft JhengHei")
    run.font.size = Pt(size)
    if color:
        run.font.color.rgb = color


def h1(doc, text):
    p = doc.add_paragraph()
    r = p.add_run(text)
    style_run(r, 14, True, RGBColor(0x2A, 0x46, 0x3C))
    p.paragraph_format.space_before = Pt(14)
    p.paragraph_format.space_after = Pt(6)


def h2(doc, text):
    p = doc.add_paragraph()
    r = p.add_run(text)
    style_run(r, 11.5, True)
    p.paragraph_format.space_before = Pt(10)
    p.paragraph_format.space_after = Pt(4)


def para(doc, text, size=10.5):
    p = doc.add_paragraph()
    r = p.add_run(text)
    style_run(r, size)
    p.paragraph_format.space_after = Pt(5)


def bullets(doc, items):
    for item in items:
        p = doc.add_paragraph(item, style="List Bullet")
        for r in p.runs:
            style_run(r, 10)


def table(doc, headers, rows, col_widths=None):
    t = doc.add_table(rows=1 + len(rows), cols=len(headers))
    t.style = "Table Grid"
    t.alignment = WD_TABLE_ALIGNMENT.CENTER
    for i, h in enumerate(headers):
        t.rows[0].cells[i].text = h
        set_cell_shading(t.rows[0].cells[i], "D5E8F0")
        for p in t.rows[0].cells[i].paragraphs:
            for r in p.runs:
                style_run(r, 9, True)
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            t.rows[ri + 1].cells[ci].text = str(val)
            for p in t.rows[ri + 1].cells[ci].paragraphs:
                for r in p.runs:
                    style_run(r, 9)
    if col_widths:
        for row in t.rows:
            for i, w in enumerate(col_widths):
                row.cells[i].width = Cm(w)
    doc.add_paragraph()
    return t


def build_docx():
    doc = Document()
    for s in doc.sections:
        s.top_margin = Cm(2)
        s.bottom_margin = Cm(2)
        s.left_margin = Cm(2.5)
        s.right_margin = Cm(2.5)

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run("養康 KOL 行銷分析\n及同級別預算參考")
    style_run(r, 20, True, RGBColor(0x2A, 0x46, 0x3C))
    para(doc, "KOL Marketing Analysis — 頤安本草對標研究", 11)
    para(
        doc,
        f"日期：{today.year} 年 {today.month} 月 {today.day} 日　｜"
        "資料來源：Instagram 公開帖文及帳號資料　｜"
        "對標：養康中醫館 @hons_cmc",
        9,
    )
    doc.add_paragraph()

    h1(doc, "一、摘要")
    para(
        doc,
        "養康中醫館長期以「KOL 到店體驗＋真實分享」作社交種草，合作對象以微型至大型 KOL 為主。"
        "內容以針灸、拔罐、艾灸、中藥調理等體驗式帖文為主，強調個人故事（肩頸痛、失眠、濕疹、腸胃等），"
        "而非硬銷優惠。以下整理 8 位曾合作 KOL，並估算 2026 年香港同級別市場報價。",
    )
    bullets(doc, [
        "合作 KOL 層級：微型 1 位｜小型 1 位（Amy Ng 1.8 萬粉）｜中階 4 位｜大型 2 位",
        "高互動帖文讚好：366 – 7,243（Hana Tam 帖文最高）",
        "頤安建議：採「中階＋小型」體驗式合作，性價比與養康主力策略一致",
        "KOL 費用不含於月度 HKD 10,000 行銷月費，建議另列年度 KOL 預算",
    ])

    h1(doc, "二、養康 KOL 合作清單")
    kol_rows = [
        (k["name"], "@" + k["handle"], k["followers"], k["tier"],
         k["treatment"], k["likes"], k["year"])
        for k in KOLS
    ]
    table(
        doc,
        ["KOL", "IG 帳號", "粉絲", "層級", "合作內容", "帖文讚好", "約年份"],
        kol_rows,
        [2.2, 2.5, 1.8, 1.8, 3.5, 1.5, 1.2],
    )
    para(doc, "IG 連結：", 10)
    for k in KOLS:
        para(doc, f"• {k['name']} — {k['url']}", 9)

    h1(doc, "三、內容模式分析")
    table(doc, ["模式", "說明", "養康案例"], [
        ("體驗式種草", "KOL 親身到店，分享治療過程與感受", "針灸落針、診所環境、農本方"),
        ("痛點共鳴", "以 OL 肩頸痛、失眠、濕疹、腸胃等切入", "meikiw 腰痛、eva 腸胃、bithialam 濕疹"),
        ("信任背書", "展示診所、醫師互動、治療細節", "stephanie 首次針灸緊張→放鬆"),
        ("標記官方", "帖文 @hons_cmc 導流", "eva_pinkland、hanatam、stephanie_ssw"),
    ], [2.5, 5, 8])

    h2(doc, "頤安可借鑑 vs 不照搬")
    bullets(doc, [
        "借鑑：體驗式真實分享、症狀故事、多圖／Reels 展示過程",
        "借鑑：中階 KOL（5–15 萬粉）性價比較高，信任度佳",
        "不照搬：「3 日濕疹好七成」等進取療效表述（合規風險）",
        "不照搬：大型 KOL 單帖 HKD 2.5 萬以上，需按預算及 ROI 選用",
        "不照搬：免診金／首診優惠換取曝光（與頤安定價策略不符）",
    ])

    h1(doc, "四、2026 年香港同級別 KOL 市場報價（參考）")
    para(
        doc,
        "以下為香港 IG 贊助合作之市場參考區間（2025–2026），實際報價因 KOL 議價、"
        "內容形式、拍攝要求、使用授權期及是否含 Reels／Story 而異。",
        10,
    )
    table(
        doc,
        ["層級", "粉絲參考", "單圖 Feed 帖 HKD", "Reels 影片 HKD", "年約 4 次合作 HKD"],
        list(PRICE_TIERS),
        [1.5, 2.5, 2.8, 2.8, 3],
    )
    h2(doc, "4.1 額外費用因素")
    bullets(doc, [
        "到店體驗：多數 KOL 要求免費或折扣診療＋交通；頂級另加 appearance fee",
        "多圖輪播／Reels 拍攝：通常比單圖貴 1.5–2.5 倍",
        "Story 限時動態：約為 Feed 帖之 30–50%",
        "內容授權（用於廣告投放）：另加 30–100%",
        "藝人／電視主持：同等粉絲下報價通常高 30–80%",
    ])
    para(
        doc,
        "備註：養康多數案例為 2015 年前後，當時 KOL 報價約為現時 1/3–1/5；"
        "以上為 2026 年市場參考，非歷史成交價。",
        9,
    )

    h1(doc, "五、頤安 KOL 預算方案（建議）")
    table(
        doc,
        ["方案", "頻率", "組合", "年度預算", "月均"],
        list(BUDGET_PLANS),
        [2, 1.5, 4.5, 2.5, 2],
    )
    para(doc, "建議頤安採方案 B（標準），理由：", 10)
    bullets(doc, [
        "與養康「中階 KOL 為主」策略相近（meikiw、cherrietam、bithialam 同級）",
        "年度 HKD 64,000 可配合 Meta 廣告將 KOL 帖文加推（spark ads）",
        "配合月度 4 IG + 4 FB 自有內容，KOL 作季度信任背書即可",
        "KOL 費用建議客戶直接支付或另列「KOL 專項預算」，不含 HKD 10,000 月費",
    ])

    h1(doc, "六、執行要點")
    bullets(doc, [
        "人選：中環上班族、媽媽、美容健康類中階 KOL；優先 5–15 萬粉絲",
        "內容：合規體驗文＋ WhatsApp 6734 9532 CTA；避免誇大療效",
        "流程：Brief → 到店體驗 → 草稿審批（5 工作天）→ 發布 → 可選 Meta 加推",
        "量度：追蹤 @oakville.wellness 引流、whatsapp_click、KOL 帖文 URL UTM",
    ])

    h1(doc, "七、合規提醒")
    bullets(doc, [
        "KOL 帖文須符合香港醫療廣告規範；療效描述由客戶及 KOL 共同把關",
        "承辦方可代審草稿，但最終發布內容責任由客戶確認",
        "不建議使用「治癒」「保證見效」「X 日好 Y 成」等表述",
    ])

    doc.save(DOCX_OUT)
    shutil.copy2(DOCX_OUT, DOCX_DESKTOP)
    os.makedirs(os.path.dirname(DOCX_PRES), exist_ok=True)
    shutil.copy2(DOCX_OUT, DOCX_PRES)


# ── PPTX helpers ──

def set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    set_fill(s, color)
    s.line.fill.background()
    return s


def add_textbox(slide, left, top, width, height, text, size=14, bold=False,
                color=INK, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    r = p.runs[0]
    r.font.name = FONT
    r.font.size = PptPt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def add_bullets(slide, left, top, width, height, items, size=11, color=INK):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.space_after = PptPt(4)
        r = p.runs[0]
        r.font.name = FONT
        r.font.size = PptPt(size)
        r.font.color.rgb = color
    return tb


def slide_header(slide, num, title, subtitle=""):
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, PAPER)
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), CINNABAR)
    add_rect(slide, 0, Inches(0.08), SLIDE_W, Inches(1.05), PINE)
    if num:
        add_textbox(slide, Inches(0.55), Inches(0.22), Inches(1.2), Inches(0.35),
                    num, 11, True, OCHRE)
    add_textbox(slide, Inches(0.55), Inches(0.48), Inches(12), Inches(0.55),
                title, 26, True, WHITE)
    if subtitle:
        add_textbox(slide, Inches(0.55), Inches(0.95), Inches(12), Inches(0.35),
                    subtitle, 11, False, PAPER)
    add_textbox(slide, Inches(0.55), Inches(7.05), Inches(6), Inches(0.3),
                "頤安本草 · KOL 行銷分析", 9, False, MUTED)


def build_pptx():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    # 1 Cover
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, PINE)
    add_rect(s, 0, 0, Inches(0.12), SLIDE_H, CINNABAR)
    add_textbox(s, Inches(0.9), Inches(1.8), Inches(8), Inches(0.45), "KOL ANALYSIS", 13, False, OCHRE)
    add_textbox(s, Inches(0.9), Inches(2.35), Inches(11), Inches(1.0),
                "養康 KOL 行銷分析\n同級別預算參考", 34, True, WHITE)
    add_textbox(s, Inches(0.9), Inches(3.55), Inches(10), Inches(0.8),
                f"頤安本草對標研究　｜　{today.year} 年 6 月\n"
                "資料：Instagram @hons_cmc 合作案例",
                14, False, PAPER)

    # 2 Overview
    s = prs.slides.add_slide(blank)
    slide_header(s, "01", "養康 KOL 策略概覽", "體驗式種草 · 症狀故事 · @hons_cmc 導流")
    add_rect(s, Inches(0.55), Inches(1.42), Inches(3.8), Inches(1.35), WHITE)
    add_textbox(s, Inches(0.75), Inches(1.55), Inches(3.4), Inches(0.65), "8", 36, True, PINE)
    add_textbox(s, Inches(0.75), Inches(2.15), Inches(3.4), Inches(0.4), "已整理合作 KOL", 11, False, MUTED)
    add_rect(s, Inches(4.55), Inches(1.42), Inches(3.8), Inches(1.35), WHITE)
    add_textbox(s, Inches(4.75), Inches(1.55), Inches(3.4), Inches(0.65), "4", 36, True, CINNABAR)
    add_textbox(s, Inches(4.75), Inches(2.15), Inches(3.4), Inches(0.4), "層級（微→大型）", 11, False, MUTED)
    add_rect(s, Inches(8.55), Inches(1.42), Inches(4.2), Inches(1.35), WHITE)
    add_textbox(s, Inches(8.75), Inches(1.55), Inches(3.8), Inches(0.65), "7,243", 32, True, PINE)
    add_textbox(s, Inches(8.75), Inches(2.15), Inches(3.8), Inches(0.4), "最高帖文讚好（Hana Tam）", 10, False, MUTED)
    add_rect(s, Inches(0.55), Inches(2.95), Inches(12.2), Inches(3.85), WHITE)
    add_bullets(s, Inches(0.75), Inches(3.1), Inches(11.8), Inches(3.5), [
        "微型：pinksheeb（約 4k 粉，拔罐艾灸）｜小型：amyng_amy（1.8 萬粉，美妝／生活）",
        "中階：meikiw、cherrietam、bithialam、stephanie_ssw（5.5 萬粉，針灸體驗）",
        "大型：eva_pinkland、hanatam（@hons_cmc，5k–7k 讚）",
        "共通：到店真實體驗＋個人健康故事，非硬銷優惠",
    ], 11.5)

    # 3 KOL table (compact)
    s = prs.slides.add_slide(blank)
    slide_header(s, "02", "KOL 合作清單")
    rows_data = [
        ("Eva Cheung", "@eva_pinkland", "26 萬", "大型", "腸胃+中藥", "5,806"),
        ("Hana Tam", "@hanatam", "40 萬", "大型", "針灸助眠", "7,243"),
        ("黃美棋", "@meikiw", "11 萬", "中階", "腰痛針灸", "1,091"),
        ("申倩瑋", "@stephanie_ssw", "5.5 萬", "中階", "針灸體驗", "3,320"),
        ("Cherrie", "@cherrietam", "5.8 萬", "中階", "養康腸胃", "678"),
        ("Bithia Lam", "@bithialam", "4.3 萬", "中階", "濕疹調理", "366"),
        ("Amy Ng", "@amyng_amy", "1.8 萬", "小型", "OL 肩頸", "9"),
        ("Sandra Ho", "@pinksheeb", "4,237", "微型", "拔罐艾灸", "9"),
    ]
    y = Inches(1.45)
    xs = [Inches(0.55), Inches(2.5), Inches(4.2), Inches(5.5), Inches(7.2), Inches(10.5)]
    ws = [Inches(1.9), Inches(1.65), Inches(1.25), Inches(1.65), Inches(3.25), Inches(2.2)]
    headers = ["KOL", "帳號", "粉絲", "層級", "內容", "讚好"]
    for i, h in enumerate(headers):
        add_rect(s, xs[i], y, ws[i], Inches(0.36), PINE)
        add_textbox(s, xs[i] + Inches(0.06), y + Inches(0.05), ws[i] - Inches(0.1), Inches(0.28),
                    h, 9, True, WHITE)
    y += Inches(0.36)
    for ri, row in enumerate(rows_data):
        bg = PAPER if ri % 2 == 0 else WHITE
        rh = Inches(0.52)
        for ci, cell in enumerate(row):
            add_rect(s, xs[ci], y, ws[ci], rh, bg)
            add_textbox(s, xs[ci] + Inches(0.06), y + Inches(0.08), ws[ci] - Inches(0.1), rh - Inches(0.1),
                        cell, 8.5, ci == 0, PINE if ci == 0 else INK)
        y += rh

    # 4 Pricing tiers
    s = prs.slides.add_slide(blank)
    slide_header(s, "03", "2026 年同級別市場報價（參考）", "單次合作 · 港元 HKD")
    y = Inches(1.48)
    for i, (tier, fans, feed, reels, annual) in enumerate(PRICE_TIERS):
        bg = PAPER if i % 2 == 0 else WHITE
        add_rect(s, Inches(0.55), y, Inches(12.2), Inches(0.72), bg)
        add_textbox(s, Inches(0.7), y + Inches(0.12), Inches(1.3), Inches(0.48), tier, 10, True, PINE)
        add_textbox(s, Inches(2.1), y + Inches(0.12), Inches(2.2), Inches(0.48), fans, 9, False, MUTED)
        add_textbox(s, Inches(4.4), y + Inches(0.12), Inches(2.5), Inches(0.48), f"Feed {feed}", 9.5, False, INK)
        add_textbox(s, Inches(7.0), y + Inches(0.12), Inches(2.8), Inches(0.48), f"Reels {reels}", 9.5, False, INK)
        add_textbox(s, Inches(9.9), y + Inches(0.12), Inches(2.6), Inches(0.48), f"年4次 {annual}", 9.5, True, CINNABAR)
        y += Inches(0.72)
    add_rect(s, Inches(0.55), Inches(5.35), Inches(12.2), Inches(1.35), PINE_LIGHT)
    add_textbox(s, Inches(0.75), Inches(5.5), Inches(11.8), Inches(1.05),
                "2015 年養康合作價約為現時 1/3–1/5。另計：Story、廣告授權、藝人加成、到店體驗成本。",
                10, False, WHITE)

    # 5 Budget plans
    s = prs.slides.add_slide(blank)
    slide_header(s, "04", "頤安 KOL 預算方案", "建議採方案 B · 不含月度 HKD 10,000 月費")
    for i, (name, freq, combo, annual, monthly) in enumerate(BUDGET_PLANS):
        col = i % 2
        row = i // 2
        x = Inches(0.55) + col * Inches(6.35)
        y = Inches(1.48) + row * Inches(2.55)
        color = PINE if i == 1 else (PINE_LIGHT if i == 0 else (OCHRE if i == 2 else CINNABAR))
        add_rect(s, x, y, Inches(6.0), Inches(2.35), WHITE)
        add_rect(s, x, y, Inches(6.0), Inches(0.42), color)
        add_textbox(s, x + Inches(0.12), y + Inches(0.05), Inches(5.7), Inches(0.32), name, 13, True, WHITE)
        add_textbox(s, x + Inches(0.12), y + Inches(0.55), Inches(5.7), Inches(1.65),
                    f"{freq}　{combo}\n\n{annual}　（{monthly}）", 10.5, False, INK)
    add_rect(s, Inches(0.55), Inches(6.55), Inches(12.2), Inches(0.55), PINE)
    add_textbox(s, Inches(0.75), Inches(6.62), Inches(11.8), Inches(0.4),
                "方案 B 約 HKD 64,000/年 ≈ 月均 5,300　｜　配合 Meta 加推 KOL 帖文效果更佳",
                10, False, PAPER)

    # 6 Conclusion
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, PINE)
    add_rect(s, 0, 0, Inches(0.12), SLIDE_H, CINNABAR)
    add_textbox(s, Inches(0.9), Inches(1.5), Inches(11), Inches(0.4), "結論", 13, False, OCHRE)
    add_bullets(s, Inches(0.9), Inches(2.05), Inches(11.5), Inches(3.5), [
        "養康以「中階 KOL 體驗文」為主力（含申倩瑋 5.5 萬粉），大型 KOL 作加強曝光",
        "2026 同級別中階 KOL 單帖約 HKD 8,000–25,000；年度 4–6 次約 HKD 24,000–64,000",
        "頤安建議：方案 B + 合規體驗文 + WhatsApp 6734 9532",
        "不照搬：誇大療效、免診金、大型 KOL 高預算無 ROI 追蹤",
        "詳見 Oakville-KOL-Marketing-Analysis.docx",
    ], 14, WHITE)

    prs.save(PPTX_OUT)
    shutil.copy2(PPTX_OUT, PPTX_DESKTOP)
    os.makedirs(os.path.dirname(PPTX_PRES), exist_ok=True)
    shutil.copy2(PPTX_OUT, PPTX_PRES)


if __name__ == "__main__":
    build_docx()
    build_pptx()
    print(f"DOCX: {DOCX_OUT}")
    print(f"DOCX Desktop: {DOCX_DESKTOP}")
    print(f"PPTX: {PPTX_OUT}")
    print(f"PPTX Desktop: {PPTX_DESKTOP}")
