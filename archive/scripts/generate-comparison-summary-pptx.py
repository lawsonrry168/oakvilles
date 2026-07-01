# -*- coding: utf-8 -*-
"""Generate PPTX from presentations/oakville-comparison-summary.html content."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import datetime
import os
import shutil

ROOT = os.path.join(r"c:\Users\Steriod\Desktop\oakvilles")
OUTPUT = os.path.join(ROOT, "Oakville-Comparison-Summary.pptx")
OUTPUT_PRESENTATIONS = os.path.join(ROOT, "presentations", "Oakville-Comparison-Summary.pptx")

PINE = RGBColor(0x2A, 0x46, 0x3C)
PINE_LIGHT = RGBColor(0x3D, 0x5A, 0x4E)
PAPER = RGBColor(0xF7, 0xF2, 0xE7)
PAPER_DARK = RGBColor(0xE8, 0xE0, 0xD0)
CINNABAR = RGBColor(0xA2, 0x3A, 0x2E)
OCHRE = RGBColor(0xAE, 0x8A, 0x47)
INK = RGBColor(0x1A, 0x1A, 0x18)
MUTED = RGBColor(0x5A, 0x5A, 0x52)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BENCH = RGBColor(0x6B, 0x8F, 0x7A)
OLD_GRAY = RGBColor(0x9A, 0x90, 0x88)

FONT = "Microsoft JhengHei"
today = datetime.date(2026, 6, 29)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
FOOTER_L = "頤安本草 · 多維度比較總結"


def set_fill(shape, color: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color: RGBColor):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    set_fill(s, color)
    s.line.fill.background()
    return s


def add_textbox(slide, left, top, width, height, text, size=14, bold=False,
                color=INK, align=PP_ALIGN.LEFT, line_spacing=1.2):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.line_spacing = line_spacing
    r = p.runs[0]
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def add_bullets(slide, left, top, width, height, items, size=12, color=INK, spacing=Pt(4)):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = spacing
        p.line_spacing = 1.1
        r = p.runs[0]
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return tb


def slide_footer(slide):
    add_textbox(slide, Inches(0.55), Inches(7.05), Inches(6), Inches(0.3), FOOTER_L, 9, False, MUTED)
    add_textbox(slide, Inches(10.5), Inches(7.05), Inches(2.3), Inches(0.3),
                today.strftime("%Y-%m-%d"), 9, False, MUTED, PP_ALIGN.RIGHT)


def slide_header(slide, section_num: str, title: str, subtitle: str = ""):
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, PAPER)
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), CINNABAR)
    add_rect(slide, 0, Inches(0.08), SLIDE_W, Inches(1.05), PINE)
    if section_num:
        add_textbox(slide, Inches(0.55), Inches(0.22), Inches(1.2), Inches(0.35),
                    section_num, 11, True, OCHRE)
    add_textbox(slide, Inches(0.55), Inches(0.48), Inches(12), Inches(0.55),
                title, 24, True, WHITE)
    if subtitle:
        add_textbox(slide, Inches(0.55), Inches(0.95), Inches(12), Inches(0.35),
                    subtitle, 10.5, False, PAPER_DARK)
    slide_footer(slide)


def stat_card(slide, left, top, number, label, w=Inches(2.85), h=Inches(1.35)):
    add_rect(slide, left, top, w, h, WHITE)
    add_rect(slide, left, top, Inches(0.06), h, CINNABAR)
    add_textbox(slide, left + Inches(0.18), top + Inches(0.1), w - Inches(0.28), Inches(0.6),
                number, 30, True, PINE)
    add_textbox(slide, left + Inches(0.18), top + Inches(0.72), w - Inches(0.28), Inches(0.5),
                label, 9.5, False, MUTED)


def comparison_col(slide, left, top, width, header, header_color, items, body_h=Inches(4.2), size=10.5):
    h = Inches(0.4)
    add_rect(slide, left, top, width, h, header_color)
    add_textbox(slide, left + Inches(0.1), top + Inches(0.05), width - Inches(0.2), h,
                header, 11.5, True, WHITE)
    add_rect(slide, left, top + h, width, body_h, WHITE)
    add_bullets(slide, left + Inches(0.1), top + h + Inches(0.1),
                width - Inches(0.2), body_h - Inches(0.12), items, size, INK, Pt(3))


def data_table(slide, y_start, headers, rows, col_widths, row_h=Inches(0.58), header_h=Inches(0.36), font_size=9):
    xs = [Inches(0.55)]
    for w in col_widths[:-1]:
        xs.append(xs[-1] + w)
    y = y_start
    for i, h in enumerate(headers):
        add_rect(slide, xs[i], y, col_widths[i], header_h, PINE)
        add_textbox(slide, xs[i] + Inches(0.06), y + Inches(0.04), col_widths[i] - Inches(0.1),
                    header_h, h, 9, True, WHITE)
    y += header_h
    for ri, row in enumerate(rows):
        bg = WHITE if ri % 2 == 0 else PAPER
        for ci, cell in enumerate(row):
            add_rect(slide, xs[ci], y, col_widths[ci], row_h, bg)
            add_textbox(slide, xs[ci] + Inches(0.06), y + Inches(0.05),
                        col_widths[ci] - Inches(0.1), row_h - Inches(0.08),
                        cell, font_size, ci == 0, PINE if ci == 0 else INK)
        y += row_h


prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]

# ── Cover ──
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, PINE)
add_rect(s, 0, 0, Inches(0.12), SLIDE_H, CINNABAR)
add_textbox(s, Inches(0.9), Inches(1.55), Inches(8), Inches(0.4),
            "OAKVILLE WELLNESS · STRATEGY MATRIX", 12, False, OCHRE)
add_textbox(s, Inches(0.9), Inches(2.05), Inches(10.5), Inches(1.05),
            "頤安本草官網重建\n多維度比較總結", 36, True, WHITE)
add_textbox(s, Inches(0.9), Inches(3.25), Inches(11), Inches(0.95),
            "整合《Oakville-Website-Proposal》方案書與簡報內容，"
            "對照 2026 年 6 月已部署新站之實際狀態。"
            "三方對比：養康（標竿）｜舊站｜新站（現況）。",
            13, False, PAPER)
legend_y = Inches(4.45)
for i, (color, label) in enumerate([
    (BENCH, "養康中醫館（對標）"),
    (OLD_GRAY, "舊站 oakvilles.com"),
    (PINE, "新站（已上線）"),
]):
    x = Inches(0.9) + i * Inches(3.8)
    dot = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, legend_y, Inches(0.18), Inches(0.18))
    set_fill(dot, color)
    dot.line.fill.background()
    add_textbox(s, x + Inches(0.28), legend_y - Inches(0.02), Inches(3.4), Inches(0.28),
                label, 11, False, PAPER_DARK)
add_textbox(s, Inches(0.9), Inches(5.15), Inches(10), Inches(0.35),
            f"{today.strftime('%Y 年 %m 月')} ｜ oakvilles.vercel.app ｜ WhatsApp 6734 9532", 11, False, PAPER_DARK)

# ── 01 一句話結論 ──
s = prs.slides.add_slide(blank)
slide_header(s, "01", "一句話結論")
add_textbox(s, Inches(0.55), Inches(1.28), Inches(12.2), Inches(0.75),
            "借鑑養康的「內容厚度 + WhatsApp 轉化」，補足舊站 SEO／信任／量度缺口；"
            "新站以中環高端、合規克制定位與其促銷密集型策略區隔。",
            12.5, False, INK)
stat_card(s, Inches(0.55), Inches(2.2), "1→62", "頁面規模（含雙語）")
stat_card(s, Inches(3.55), Inches(2.2), "0→6", "可追蹤轉化事件")
stat_card(s, Inches(6.55), Inches(2.2), "7", "症狀 SEO 著陸頁")
stat_card(s, Inches(9.55), Inches(2.2), "5.0", "Google 評價展示")

# ── 02 核心維度對照（上）──
s = prs.slides.add_slide(blank)
slide_header(s, "02", "核心維度對照表（上）", "養康 honscmc.com ｜ 舊站 ｜ 新站（現況）")
data_table(s, Inches(1.38),
           ["維度", "養康", "舊站", "新站（現況）"],
           [
               ("定位", "尖沙咀·皮膚/美容·促銷", "個人醫師·訊息分散", "中環高端·頤安本草·合規克制"),
               ("頁面規模", "WP 多頁·四層 nav", "實質 1 頁·nav 5 項", "62 頁（中+英）·已部署 Vercel"),
               ("SEO 入口", "每症狀一頁·地區+症狀", "無症狀頁·無 sitemap", "7 症狀+central-hk·Schema·OG"),
               ("轉化動線", "Joinchat·首診優惠", "僅 WA 浮窗", "2 步 funnel·計算器·sticky CTA"),
           ],
           [Inches(1.35), Inches(3.35), Inches(3.2), Inches(3.85)],
           row_h=Inches(0.72), font_size=9.5)

# ── 02b 核心維度對照（下）──
s = prs.slides.add_slide(blank)
slide_header(s, "02", "核心維度對照表（下）")
data_table(s, Inches(1.38),
           ["維度", "養康", "舊站", "新站（現況）"],
           [
               ("信任元件", "13年·8000診量·FAQ", "缺評價·環境·FAQ", "Google 5.0·gallery·明碼價"),
               ("視覺風格", "綠色臨床·密度高", "通用 SPA·Hero 輪播", "paper/pine 禪意·Mobile-first"),
               ("量度能力", "（未公開 GA）", "無 GA4·無 CTA", "6 事件就緒·GA4 待填"),
               ("合規風險", "efficacy 較進取", "數週見效·文案筆誤", "禁止誇大·評價免責·低風險"),
           ],
           [Inches(1.35), Inches(3.35), Inches(3.2), Inches(3.85)],
           row_h=Inches(0.72), font_size=9.5)

# ── 03 借鑑 vs 不照搬 ──
s = prs.slides.add_slide(blank)
slide_header(s, "03", "從對標「借鑑 vs 不照搬」")
comparison_col(s, Inches(0.55), Inches(1.38), Inches(3.95), "✓ 借鑑養康", BENCH, [
    "以「症狀」而非抽象科別作 SEO 入口",
    "FAQ 結構化（地點、資歷、收費、預約）",
    "WhatsApp 全站一致預約動線",
    "Blog 持續內容 + 內部連結",
    "首屏展示地點、時間、專科範圍",
], Inches(4.85))
comparison_col(s, Inches(4.7), Inches(1.38), Inches(3.95), "✗ 不照搬", CINNABAR, [
    "美容塑形促銷語（豐胸、減肥主入口）",
    "首診免診金等高折扣（高端定位不符）",
    "efficacy 進取表述 · 論壇軟文大量佈局",
    "資訊過密、視覺擁擠的臨床促銷感",
], Inches(4.85))
comparison_col(s, Inches(8.85), Inches(1.38), Inches(3.9), "頤安差異化", PINE, [
    "養康：流量 + 轉化優先",
    "頤安：信任 + 合規 + 高端體驗優先",
    "架構可借鑑，品牌調性刻意區隔",
], Inches(4.85))

# ── 04 12 項改善（上）──
s = prs.slides.add_slide(blank)
slide_header(s, "04", "舊站 → 新站：12 項改善（上）")
data_table(s, Inches(1.38),
           ["類別", "舊站", "新站"],
           [
               ("頁面", "1 頁為主", "62 頁（含 /en/ 雙語）"),
               ("品牌", "伍厚臻個人名義", "頤安本草 + 印章識別"),
               ("SEO", "基本 title", "sitemap·JSON-LD·canonical·OG 1200×630"),
               ("本地", "無中環專頁", "about/central-hk · GEO · llms.txt"),
               ("預約", "WA 連結", "2 步 funnel + sessionStorage 串接"),
               ("收費", "未展示", "明碼價 + 即時計算器"),
           ],
           [Inches(1.55), Inches(5.2), Inches(5.4)],
           row_h=Inches(0.62), font_size=10)

# ── 04b 12 項改善（下）──
s = prs.slides.add_slide(blank)
slide_header(s, "04", "舊站 → 新站：12 項改善（下）")
data_table(s, Inches(1.38),
           ["類別", "舊站", "新站"],
           [
               ("評價", "無", "Google 5.0 + Review Schema"),
               ("量度", "無", "whatsapp_click 等 6 事件"),
               ("內容", "2 專科卡片", "4 專科 + 7 症狀 + Blog + FAQ"),
               ("合規", "有效能承諾", "禁止誇大表述"),
               ("手機", "基本 RWD", "sticky CTA · WA 浮動避障"),
               ("部署", "舊 hosting", "GitHub → Vercel 自動部署"),
           ],
           [Inches(1.55), Inches(5.2), Inches(5.4)],
           row_h=Inches(0.62), font_size=10)

# ── 05 行銷策略 ──
s = prs.slides.add_slide(blank)
slide_header(s, "05", "行銷與廣告策略", "提案 + 現站就緒度")
data_table(s, Inches(1.38),
           ["策略", "養康做法", "頤安建議", "新站準備度"],
           [
               ("SEO 長尾", "症狀+地區著陸頁", "持續補 FAQ／症狀頁", "7 頁已上線"),
               ("WhatsApp", "24h · Joinchat", "統一 WA + 全站追蹤", "funnel 已就緒"),
               ("Google 本地", "—", "GBP 優化 + 評價邀請", "待執行"),
               ("Meta 廣告", "療程銷售頁", "症狀頁著陸·Pixel 追蹤", "Pixel 待接"),
               ("Google Ads", "—", "品牌+症狀+中環 Search", "GA4 ID 待填"),
               ("IG 種草", "KOL 素人化", "@oakville.wellness 非硬銷", "官網互導已接"),
               ("首診優惠", "皮膚科免診金", "不建議（高端定位）", "—"),
               ("內容節奏", "Blog 持續更新", "每月 1–2 篇", "4 篇已上線"),
           ],
           [Inches(1.45), Inches(2.85), Inches(3.55), Inches(3.9)],
           row_h=Inches(0.52), font_size=8.5)

# ── 06 廣告著陸頁 ──
s = prs.slides.add_slide(blank)
slide_header(s, "06", "廣告著陸頁對應", "Meta / Google 共用")
data_table(s, Inches(1.55),
           ["廣告主題", "建議著陸頁", "主要 KPI"],
           [
               ("濕疹 / 暗瘡", "/conditions/eczema.html · /conditions/acne.html", "whatsapp_click"),
               ("備孕 / 婦科", "/conditions/fertility.html", "whatsapp_click"),
               ("頸痛 / 痛症", "/conditions/neck-pain.html · /services/pain.html", "whatsapp_click"),
               ("中環中醫", "/about/central-hk.html", "cta_click · booking_submit"),
               ("外籍客群", "/en/ · /en/services/skin.html", "whatsapp_click"),
               ("品牌詞", "/ · /about.html", "cta_click"),
           ],
           [Inches(2.2), Inches(6.5), Inches(3.6)],
           row_h=Inches(0.68), font_size=10)

# ── 07 交付路線圖 ──
s = prs.slides.add_slide(blank)
slide_header(s, "07", "交付路線圖", "提案 vs 2026.06 現況")
data_table(s, Inches(1.48),
           ["階段", "提案內容", "2026.06 現況"],
           [
               ("第一階段", "30 頁 · SEO · funnel · GA4 就緒", "✓ 已完成並超越（62 頁雙語）"),
               ("第二階段", "50 張品牌視覺資產全站接入", "規劃完成 · 部分已接入"),
               ("第三階段", "Production 上線 oakvilles.com", "✓ 已部署 Vercel"),
               ("持續", "月度行銷 HKD 10,000", "可選 · 待啟動"),
               ("廣告開跑前", "—", "GA4·GTM·Pixel·Search Console·GBP"),
           ],
           [Inches(1.65), Inches(5.1), Inches(5.4)],
           row_h=Inches(0.72), font_size=10)

# ── 08 建議下一步 ──
s = prs.slides.add_slide(blank)
slide_header(s, "08", "建議下一步（優先序）")
cols = [
    ("P0 · 本週", PINE, [
        "填入 SITE_GA4_ID 並部署",
        "Search Console + sitemap",
        "Google Business Profile 優化",
        "oakvilles.com DNS 指向 Vercel",
    ]),
    ("P1 · 廣告前", OCHRE, [
        "GTM：GA4 + Meta Pixel",
        "GA4 轉化 → Google Ads",
        "Meta 自訂轉換 whatsapp_click",
    ]),
    ("P2 · 持續", CINNABAR, [
        "症狀 campaign + 著陸 A/B",
        "Blog 每月 1–2 篇",
        "品牌視覺資產批次接入",
    ]),
]
for i, (title, color, items) in enumerate(cols):
    x = Inches(0.55) + i * Inches(4.15)
    add_rect(s, x, Inches(1.42), Inches(3.85), Inches(5.15), WHITE)
    add_rect(s, x, Inches(1.42), Inches(3.85), Inches(0.4), color)
    add_textbox(s, x + Inches(0.12), Inches(1.48), Inches(3.6), Inches(0.3), title, 13, True, WHITE)
    add_bullets(s, x + Inches(0.12), Inches(2.0), Inches(3.6), Inches(4.3), items, 11.5, INK)

# ── Closing ──
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, PINE)
add_rect(s, 0, 0, Inches(0.12), SLIDE_H, CINNABAR)
add_textbox(s, Inches(0.9), Inches(1.65), Inches(11), Inches(0.35), "資料來源", 12, False, OCHRE)
add_textbox(s, Inches(0.9), Inches(2.15), Inches(11.2), Inches(1.1),
            "Oakville-Website-Proposal.docx · Oakville-Website-Proposal.pptx（2026-06-20）\n"
            "新站 repo 部署紀錄 · 對標 honscmc.com",
            12, False, PAPER_DARK)
add_rect(s, Inches(0.9), Inches(3.55), Inches(11.5), Inches(2.35), PINE_LIGHT)
add_textbox(s, Inches(1.1), Inches(3.75), Inches(11), Inches(0.35), "頤安本草 · 伍厚臻中醫師", 14, True, OCHRE)
add_textbox(s, Inches(1.1), Inches(4.25), Inches(11), Inches(1.45),
            "oakvilles.vercel.app ｜ WhatsApp +852 6734 9532 ｜ 電話 2881 8182\n"
            "中環皇后大道中 99 號錦安大廈 6 樓 602 室",
            13, False, PAPER)

prs.save(OUTPUT)
os.makedirs(os.path.dirname(OUTPUT_PRESENTATIONS), exist_ok=True)
shutil.copy2(OUTPUT, OUTPUT_PRESENTATIONS)
print(f"Saved: {OUTPUT}")
print(f"Copied: {OUTPUT_PRESENTATIONS}")
print(f"Slides: {len(prs.slides)}")
