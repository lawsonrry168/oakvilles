# -*- coding: utf-8 -*-
"""Generate North Star + GTM + Battlecard DOCX and PPTX for 頤安本草."""

from docx import Document
from docx.shared import Pt, Cm, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
from pptx import Presentation
from pptx.util import Inches, Pt as PptPt
from pptx.dml.color import RGBColor as PptRGB
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import datetime
import os
import shutil

BASE = r"c:\Users\Steriod\Desktop\oakvilles"
DOCX_OUT = os.path.join(BASE, "Oakville-Strategy-Pack.docx")
DOCX_DESKTOP = r"c:\Users\Steriod\Desktop\Oakville-Strategy-Pack.docx"
DOCX_PRES = os.path.join(BASE, "presentations", "Oakville-Strategy-Pack.docx")
PPTX_OUT = os.path.join(BASE, "Oakville-Strategy-Pack.pptx")
PPTX_DESKTOP = r"c:\Users\Steriod\Desktop\Oakville-Strategy-Pack.pptx"
PPTX_PRES = os.path.join(BASE, "presentations", "Oakville-Strategy-Pack.pptx")

today = datetime.date(2026, 6, 22)

PINE = PptRGB(0x2A, 0x46, 0x3C)
PAPER = PptRGB(0xF7, 0xF2, 0xE7)
CINNABAR = PptRGB(0xA2, 0x3A, 0x2E)
OCHRE = PptRGB(0xAE, 0x8A, 0x47)
INK = PptRGB(0x1A, 0x1A, 0x18)
MUTED = PptRGB(0x5A, 0x5A, 0x52)
WHITE = PptRGB(0xFF, 0xFF, 0xFF)
PINE_LIGHT = PptRGB(0x3D, 0x5A, 0x4E)
FONT = "Microsoft JhengHei"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_cell_shading(cell, fill: str):
    shading = OxmlElement("w:shd")
    shading.set(qn("w:fill"), fill)
    shading.set(qn("w:val"), "clear")
    cell._tc.get_or_add_tcPr().append(shading)


def style_run(run, size=10.5, bold=False, color=None):
    run.bold = bold
    run.font.name = "Arial"
    run._element.rPr.rFonts.set(qn("w:eastAsia"), "Microsoft JhengHei")
    run.font.size = Pt(size)
    if color:
        run.font.color.rgb = color


def h1(doc, text):
    p = doc.add_paragraph()
    r = p.add_run(text)
    style_run(r, 14, True, RGBColor(0x2A, 0x46, 0x3C))
    p.paragraph_format.space_before = Pt(14)
    p.paragraph_format.space_after = Pt(6)


def h2(doc, text):
    p = doc.add_paragraph()
    r = p.add_run(text)
    style_run(r, 11.5, True)
    p.paragraph_format.space_before = Pt(10)
    p.paragraph_format.space_after = Pt(4)


def para(doc, text, size=10.5):
    p = doc.add_paragraph()
    r = p.add_run(text)
    style_run(r, size)
    p.paragraph_format.space_after = Pt(5)


def bullets(doc, items):
    for item in items:
        p = doc.add_paragraph(item, style="List Bullet")
        for r in p.runs:
            style_run(r, 10)


def table(doc, headers, rows, col_widths=None):
    t = doc.add_table(rows=1 + len(rows), cols=len(headers))
    t.style = "Table Grid"
    t.alignment = WD_TABLE_ALIGNMENT.CENTER
    for i, h in enumerate(headers):
        t.rows[0].cells[i].text = h
        set_cell_shading(t.rows[0].cells[i], "D5E8F0")
        for p in t.rows[0].cells[i].paragraphs:
            for r in p.runs:
                style_run(r, 9, True)
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            t.rows[ri + 1].cells[ci].text = str(val)
            for p in t.rows[ri + 1].cells[ci].paragraphs:
                for r in p.runs:
                    style_run(r, 9)
    if col_widths:
        for row in t.rows:
            for i, w in enumerate(col_widths):
                row.cells[i].width = Cm(w)
    doc.add_paragraph()
    return t


def build_docx():
    doc = Document()
    for s in doc.sections:
        s.top_margin = Cm(2)
        s.bottom_margin = Cm(2)
        s.left_margin = Cm(2.5)
        s.right_margin = Cm(2.5)

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run("頤安本草 · 策略文件包")
    style_run(r, 20, True, RGBColor(0x2A, 0x46, 0x3C))
    para(doc, "North Star 量度計劃 ｜ GTM 一頁摘要 ｜ 競品 Battlecard", 12)
    para(doc, f"日期：{today.year} 年 {today.month} 月 {today.day} 日　｜　oakvilles.com", 9)
    doc.add_paragraph()

    # ═══ PART 1: North Star ═══
    h1(doc, "一、North Star 量度計劃（配合 GA4）")

    h2(doc, "1.1 業務類型")
    para(
        doc,
        "頤安本草屬於 **Transaction Game（交易型）**：客戶透過官網了解服務後，以 WhatsApp 預約診症。"
        "North Star 應反映「客戶成功取得專業中醫協助」的意圖，而非純瀏覽量。",
    )

    h2(doc, "1.2 North Star Metric")
    table(doc, ["項目", "定義"], [
        ("North Star", "whatsapp_click — 用戶點擊 WhatsApp 6734 9532 連結"),
        ("為何選它", "直接反映「想預約／查詢」的高意圖行為；全站統一 KPI；可追蹤來源與著陸頁"),
        ("非 NSM", "營收、診症人次（滯後）、Pageviews（非客戶價值）、電話來電（難追蹤）"),
    ], [3.5, 12.5])

    h2(doc, "1.3 Input Metrics（指標星座，5 項）")
    table(doc, ["#", "Input Metric", "定義", "優化槓桿"], [
        ("I1", "booking_submit", "2 步預約表單提交", "簡化 funnel、預填計算器方案"),
        ("I2", "calculator_to_booking", "診金計算器 → 捲動至預約", "計算器位置、CTA 文案"),
        ("I3", "症狀頁 Sessions", "/conditions/* 自然＋付費流量", "SEO、Meta/Google 著陸頁"),
        ("I4", "funnel_step 完成率", "Step1→Step2 轉化", "表單 UX、科別預選"),
        ("I5", "高意圖 cta_click", "hero / sticky / wa-float 點擊", "CTA 位置、手機 sticky 避障"),
    ], [0.8, 3, 5.5, 5.2])

    h2(doc, "1.4 網站已就緒事件 → GA4 對照")
    table(doc, ["dataLayer 事件", "觸發時機", "建議 GA4 參數", "轉化"], [
        ("whatsapp_click", "點擊 wa.me 連結", "cta_id, page_path, link_url", "✓ 主要"),
        ("booking_submit", "預約表單 submit", "form_id, page_path", "✓ 次要"),
        ("cta_click", "data-cta-id 按鈕", "cta_id, element, page_path", "監測"),
        ("funnel_step", "預約第 1/2 步", "step, specialty, page_path", "監測"),
        ("calculator_update", "計算器變更", "treatment, meds_days, total", "監測"),
        ("calculator_to_booking", "以此方案預約", "total, page_path", "✓ 次要"),
    ], [3, 4, 4.5, 1.5])

    h2(doc, "1.5 GA4 + GTM 部署清單（P0）")
    bullets(doc, [
        "建立 GA4 資源 → 取得 Measurement ID → 填入 js/site-config.js 的 SITE_GA4_ID",
        "GTM 容器：載入 GA4 Config + 6 個 Custom Event（由 dataLayer 觸發）",
        "GA4 標記 whatsapp_click、booking_submit 為「轉化」",
        "匯入 Google Ads 轉化；Meta Pixel 自訂轉換對齊 whatsapp_click",
        "Search Console 提交 sitemap；UTM 規範：utm_source / medium / campaign",
        "月報：NSM 趨勢 + 5 個 Input Metrics + 來源／著陸頁分解",
    ])

    h2(doc, "1.6 目標參考（上線後 90 日，待基線建立後調整）")
    table(doc, ["指標", "方向", "備註"], [
        ("whatsapp_click / 月", "↑ 月增 15–25%", "含自然 + 付費 + 社交"),
        ("WA / Sessions 率", "↑ 目標 3–5%", "全站平均；症狀頁可更高"),
        ("booking_submit / WA", "↑ 目標 20–30%", "表單 vs 直接點 WA"),
        ("症狀頁停留", "> 1 分鐘", "內容品質指標"),
        ("funnel Step1→2", "> 60%", "表單摩擦指標"),
    ], [4, 3, 9])

    # ═══ PART 2: GTM ═══
    h1(doc, "二、GTM 一頁摘要（整合行銷方案）")

    h2(doc, "2.1 產品與首攻市場")
    table(doc, ["項目", "內容"], [
        ("產品", "頤安本草 · 伍厚臻中醫師 — 中環高端全科中醫診所"),
        ("Beachhead", "中環／金鐘／半山上班族 + 外籍人士（/en/）"),
        ("核心痛點", "濕疹、暗瘡、失眠、備孕、頸肩痛 — 想搵可靠中醫"),
        ("轉化終點", "WhatsApp 6734 9532 預約"),
        ("差異化", "循本溯源、合規克制、一人一方、中環高端體驗"),
    ], [3.5, 12.5])

    h2(doc, "2.2 渠道策略")
    table(doc, ["渠道", "角色", "每月產出／預算"], [
        ("Instagram", "種草 + 信任", "4 post／月 · @oakville.wellness"),
        ("Facebook", "觸及 + 再營銷", "4 post／月（可共用 IG 素材）"),
        ("官網 Blog", "SEO 長尾", "2 篇文章／月"),
        ("Meta 廣告", "症狀受眾轉化", "HKD 5,000–8,000／月（客戶直付）"),
        ("Google Search", "高意向搜尋", "HKD 5,000–8,000／月（客戶直付）"),
        ("Google Business", "本地 Pack", "每週 1–2 則動態 + 評價邀請"),
        ("KOL（可選）", "季度信任背書", "方案 B ~64,000／年"),
    ], [2.5, 3.5, 10])

    h2(doc, "2.3 訊息主軸")
    bullets(doc, [
        "主訊息：中環 · 伍厚臻中醫師 · 循本溯源調理體質",
        "副訊息：症狀著陸頁承接（濕疹／暗瘡／備孕／失眠／痛症）",
        "CTA：WhatsApp 6734 9532 查詢檔期",
        "不採用：免診金、論壇軟文、豐胸減肥主 SEO",
    ])

    h2(doc, "2.4 成功指標（GTM KPI）")
    table(doc, ["階段", "指標", "目標"], [
        ("認知", "IG/FB 觸及、GBP 曝光", "月增"),
        ("互動", "社交互動率、症狀頁停留", "IG >2%、停留 >1min"),
        ("轉化", "whatsapp_click（NSM）", "月增 15–25%"),
        ("效率", "廣告 CPA（單次 WA 點擊）", "試跑 4–6 週後設定"),
    ], [2, 5, 9])

    h2(doc, "2.5 90 日執行路線圖")
    table(doc, ["階段", "週次", "重點"], [
        ("P0 基礎", "W1–2", "GA4 + GTM + Pixel · Search Console · GBP 優化"),
        ("試跑", "W3–6", "Meta + Google 各 1 活動 · 濕疹/備孕著陸 · 首 2 篇 Blog"),
        ("優化", "W7–10", "加碼最佳 ad set · KOL 或社交 best post 加推"),
        ("穩定", "W11–12", "月度 ROI 檢討 · 調整預算與內容日曆"),
    ], [2, 1.5, 12.5])

    # ═══ PART 3: Battlecard ═══
    h1(doc, "三、競品 Battlecard：頤安 vs 養康")

    h2(doc, "3.1 對手概覽")
    table(doc, ["項目", "養康中醫館", "頤安本草"], [
        ("網站", "honscmc.com", "oakvilles.com"),
        ("定位", "尖沙咀 · 皮膚/美容 · 促銷臨床感", "中環 · 高端全科 · 循本溯源"),
        ("品牌", "養康中醫館", "頤安本草 · 伍厚臻（003769）"),
        ("社交", "IG 活躍、KOL 種草", "@oakville.wellness 起步"),
        ("獲客", "SEO + WA + 首診優惠 + KOL", "SEO + WA + 合規內容 + 廣告"),
    ], [3, 6.5, 6.5])

    h2(doc, "3.2 能力對照")
    table(doc, ["維度", "頤安（我們）", "養康（對手）", "勝方"], [
        ("網站架構", "62 頁雙語、症狀 SEO、Schema", "WP 多頁、四層 nav", "頤安"),
        ("品牌調性", "paper/pine 禪意高端", "綠色促銷臨床感", "頤安（高端客群）"),
        ("轉化追蹤", "6 dataLayer 事件、GA4 就緒", "Joinchat、較少公開量度", "頤安"),
        ("症狀 SEO", "7 症狀頁 + central-hk", "每症狀一頁 + 地區詞", "平手"),
        ("社交/KOL", "起步中", "8+ KOL、體驗式種草", "養康"),
        ("促銷力度", "無免診金、合規克制", "首診優惠、免診金", "養康（價格敏感客）"),
        ("地區", "中環核心", "尖沙咀", "各據一方"),
        ("英文市場", "/en/ 全站", "有限", "頤安"),
        ("Google 評價", "5.0 展示", "有評價機制", "平手"),
    ], [2.2, 4.5, 4.5, 1.3])

    h2(doc, "3.3 我們贏在哪裡")
    bullets(doc, [
        "中環高端定位：吸引上班族、外籍人士，非價格敏感客群",
        "合規表述 + 一人一方：適合重視專業與信任的受眾",
        "雙語官網 + 現代 UX：預約 funnel、計算器、全站量度",
        "伍醫師個人品牌：25 年臨床、註冊 003769，差異於連鎖感",
    ])

    h2(doc, "3.4 對手贏在哪裡（及我們應對）")
    table(doc, ["養康優勢", "頤安應對"], [
        ("KOL 體驗式種草成熟", "中階 KOL 季度合作（年 ~64K）；自有 4+4 post"),
        ("首診優惠降門檻", "不跟進；強調專業與體驗價值、明碼價"),
        ("內容厚度（Blog、FAQ）", "每月 2 篇 Blog + 症狀頁持續擴充"),
        ("促銷感吸引流量", "以症狀 SEO + 廣告著陸承接，非促銷語"),
    ], [5, 11])

    h2(doc, "3.5 常見異議與回應")
    table(doc, ["客戶／內部說", "回應"], [
        ("養康有免診金", "頤安定位高端，明碼價 + 專業體驗；免診金吸引價格敏感客，非目標客群"),
        ("養康 KOL 多", "我們借鑑體驗式內容，用中階 KOL + 自有社交，合規且 ROI 可控"),
        ("養康 SEO 更久", "新站 62 頁已上線，症狀頁 + central-hk 已就位，差在時間與外鏈"),
        ("為何不做美容塑形", "與「循本溯源」定位衝突；專注全科調理與常見症狀"),
    ], [4.5, 11.5])

    h2(doc, "3.6 贏單／輸單模式")
    bullets(doc, [
        "頤安傾向贏：中環上班族、重視合規與專業、外籍、症狀搜尋高意向",
        "頤安傾向輸：價格極敏感、尖沙咀便利優先、美容塑形需求",
        "關鍵差異化：信任 + 合規 + 中環高端體驗，非價格戰",
    ])

    h2(doc, "3.7 可植入的「地雷」問題（銷售／內部用）")
    bullets(doc, [
        "「您更重視一次過便宜，還是長期調理體質？」",
        "「診所是否方便您上班途中到訪？」（中環優勢）",
        "「醫師是否親自問診、處方？」（伍醫師一人一方）",
        "「網上資訊是否清晰、預約是否方便？」（官網 + WA funnel）",
    ])

    doc.save(DOCX_OUT)
    shutil.copy2(DOCX_OUT, DOCX_DESKTOP)
    os.makedirs(os.path.dirname(DOCX_PRES), exist_ok=True)
    shutil.copy2(DOCX_OUT, DOCX_PRES)


# ── PPTX ──

def set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    set_fill(s, color)
    s.line.fill.background()
    return s


def add_textbox(slide, left, top, width, height, text, size=14, bold=False, color=INK, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    r = p.runs[0]
    r.font.name = FONT
    r.font.size = PptPt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def add_bullets(slide, left, top, width, height, items, size=11, color=INK):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.space_after = PptPt(4)
        r = p.runs[0]
        r.font.name = FONT
        r.font.size = PptPt(size)
        r.font.color.rgb = color
    return tb


def slide_header(slide, num, title, subtitle=""):
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, PAPER)
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), CINNABAR)
    add_rect(slide, 0, Inches(0.08), SLIDE_W, Inches(1.05), PINE)
    if num:
        add_textbox(slide, Inches(0.55), Inches(0.22), Inches(1.2), Inches(0.35), num, 11, True, OCHRE)
    add_textbox(slide, Inches(0.55), Inches(0.48), Inches(12), Inches(0.55), title, 26, True, WHITE)
    if subtitle:
        add_textbox(slide, Inches(0.55), Inches(0.95), Inches(12), Inches(0.35), subtitle, 11, False, PAPER)
    add_textbox(slide, Inches(0.55), Inches(7.05), Inches(6), Inches(0.3), "頤安本草 · 策略文件包", 9, False, MUTED)


def build_pptx():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    # Cover
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, PINE)
    add_rect(s, 0, 0, Inches(0.12), SLIDE_H, CINNABAR)
    add_textbox(s, Inches(0.9), Inches(2.0), Inches(11), Inches(1.2),
                "頤安本草 · 策略文件包", 36, True, WHITE)
    add_textbox(s, Inches(0.9), Inches(3.35), Inches(10), Inches(1.0),
                "North Star 量度 ｜ GTM 摘要 ｜ 競品 Battlecard\n"
                f"{today.year} 年 6 月 ｜ oakvilles.com",
                15, False, PAPER)

    # North Star
    s = prs.slides.add_slide(blank)
    slide_header(s, "01", "North Star 量度計劃", "Transaction Game · GA4 就緒")
    add_rect(s, Inches(0.55), Inches(1.42), Inches(12.2), Inches(1.15), PINE)
    add_textbox(s, Inches(0.75), Inches(1.55), Inches(11.8), Inches(0.9),
                "North Star = whatsapp_click\n"
                "用戶點擊 WhatsApp 6734 9532 — 反映「想預約／查詢」的高意圖",
                14, False, WHITE)
    add_rect(s, Inches(0.55), Inches(2.75), Inches(5.9), Inches(4.0), WHITE)
    add_textbox(s, Inches(0.75), Inches(2.88), Inches(5.5), Inches(0.3), "5 個 Input Metrics", 12, True, PINE)
    add_bullets(s, Inches(0.75), Inches(3.25), Inches(5.5), Inches(3.3), [
        "I1 booking_submit — 表單提交",
        "I2 calculator_to_booking — 計算器導預約",
        "I3 症狀頁 Sessions — SEO/廣告",
        "I4 funnel_step 完成率 — 表單 UX",
        "I5 高意圖 cta_click — hero/sticky",
    ], 10.5)
    add_rect(s, Inches(6.85), Inches(2.75), Inches(5.9), Inches(4.0), WHITE)
    add_textbox(s, Inches(7.05), Inches(2.88), Inches(5.5), Inches(0.3), "P0 部署", 12, True, CINNABAR)
    add_bullets(s, Inches(7.05), Inches(3.25), Inches(5.5), Inches(3.3), [
        "填入 SITE_GA4_ID（site-config.js）",
        "GTM：6 事件 + GA4 轉化標記",
        "Google Ads / Meta Pixel 對齊 WA",
        "月報：NSM + Input + 來源分解",
    ], 10.5)

    # GTM
    s = prs.slides.add_slide(blank)
    slide_header(s, "02", "GTM 一頁摘要", "整合月度行銷方案")
    rows = [
        ("Beachhead", "中環上班族 + 外籍 /en/"),
        ("自有內容", "4 IG + 4 FB + 2 Blog／月"),
        ("付費", "Meta 5–8K + Google 5–8K HKD／月"),
        ("訊息", "循本溯源 · 症狀著陸 · WA 6734 9532"),
        ("NSM", "whatsapp_click 月增 15–25%"),
        ("P0", "GA4 · GTM · Pixel · GBP · W1–2"),
    ]
    y = Inches(1.48)
    for i, (a, b) in enumerate(rows):
        bg = PAPER if i % 2 == 0 else WHITE
        add_rect(s, Inches(0.55), y, Inches(12.2), Inches(0.72), bg)
        add_textbox(s, Inches(0.7), y + Inches(0.12), Inches(2.5), Inches(0.48), a, 10.5, True, PINE)
        add_textbox(s, Inches(3.3), y + Inches(0.12), Inches(9.2), Inches(0.48), b, 10, False, INK)
        y += Inches(0.72)
    add_rect(s, Inches(0.55), Inches(5.95), Inches(12.2), Inches(1.0), PINE_LIGHT)
    add_textbox(s, Inches(0.75), Inches(6.1), Inches(11.8), Inches(0.75),
                "90 日：P0 基礎 → 廣告試跑 W3–6 → 優化 W7–10 → 月度 ROI W11–12",
                10, False, WHITE)

    # Battlecard
    s = prs.slides.add_slide(blank)
    slide_header(s, "03", "競品 Battlecard", "頤安 vs 養康 honscmc.com")
    add_rect(s, Inches(0.55), Inches(1.42), Inches(5.9), Inches(5.3), WHITE)
    add_textbox(s, Inches(0.75), Inches(1.55), Inches(5.5), Inches(0.3), "頤安贏", 13, True, PINE)
    add_bullets(s, Inches(0.75), Inches(1.95), Inches(5.5), Inches(4.5), [
        "中環高端 + 雙語 62 頁",
        "合規克制、一人一方",
        "GA4 6 事件量度就緒",
        "現代 UX：funnel + 計算器",
        "Google 5.0 信任背書",
    ], 10.5)
    add_rect(s, Inches(6.85), Inches(1.42), Inches(5.9), Inches(5.3), WHITE)
    add_textbox(s, Inches(7.05), Inches(1.55), Inches(5.5), Inches(0.3), "養康贏 · 我們應對", 13, True, CINNABAR)
    add_bullets(s, Inches(7.05), Inches(1.95), Inches(5.5), Inches(4.5), [
        "KOL 種草 → 中階 KOL 年 ~64K",
        "首診優惠 → 不跟；強調專業價值",
        "內容厚度 → 每月 2 Blog",
        "促銷流量 → 症狀 SEO + 廣告著陸",
    ], 10.5)
    add_rect(s, Inches(0.55), Inches(6.85), Inches(12.2), Inches(0.45), PINE)
    add_textbox(s, Inches(0.75), Inches(6.92), Inches(11.8), Inches(0.32),
                "贏單：中環上班族、重視合規專業、症狀搜尋高意向　｜　輸單：價格極敏感、尖沙咀便利、美容塑形",
                9, False, PAPER)

    prs.save(PPTX_OUT)
    shutil.copy2(PPTX_OUT, PPTX_DESKTOP)
    os.makedirs(os.path.dirname(PPTX_PRES), exist_ok=True)
    shutil.copy2(PPTX_OUT, PPTX_PRES)


if __name__ == "__main__":
    build_docx()
    build_pptx()
    print(f"DOCX: {DOCX_OUT}")
    print(f"PPTX: {PPTX_OUT}")
