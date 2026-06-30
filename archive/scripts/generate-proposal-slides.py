# -*- coding: utf-8 -*-
"""Generate integrated proposal deck — editable PPTX (python-pptx)."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import datetime
import os
import shutil

OUTPUT = os.path.join(r"c:\Users\Steriod\Desktop\oakvilles", "Oakville-Website-Proposal.pptx")
OUTPUT_DESKTOP = r"c:\Users\Steriod\Desktop\Oakville-Website-Proposal.pptx"
OUTPUT_PRESENTATIONS = os.path.join(
    r"c:\Users\Steriod\Desktop\oakvilles\presentations", "Oakville-Website-Proposal.pptx"
)

PINE = RGBColor(0x2A, 0x46, 0x3C)
PINE_LIGHT = RGBColor(0x3D, 0x5A, 0x4E)
PAPER = RGBColor(0xF7, 0xF2, 0xE7)
PAPER_DARK = RGBColor(0xE8, 0xE0, 0xD0)
CINNABAR = RGBColor(0xA2, 0x3A, 0x2E)
OCHRE = RGBColor(0xAE, 0x8A, 0x47)
INK = RGBColor(0x1A, 0x1A, 0x18)
MUTED = RGBColor(0x5A, 0x5A, 0x52)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BENCH = RGBColor(0x6B, 0x8F, 0x7A)

FONT = "Microsoft JhengHei"
today = datetime.date(2026, 6, 20)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_fill(shape, color: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color: RGBColor):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    set_fill(s, color)
    s.line.fill.background()
    return s


def add_textbox(slide, left, top, width, height, text, size=14, bold=False,
                color=INK, align=PP_ALIGN.LEFT, line_spacing=1.2):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.line_spacing = line_spacing
    r = p.runs[0]
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def add_bullets(slide, left, top, width, height, items, size=12, color=INK, spacing=Pt(5)):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = spacing
        p.line_spacing = 1.12
        r = p.runs[0]
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return tb


def slide_header(slide, section_num: str, title: str, subtitle: str = ""):
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, PAPER)
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), CINNABAR)
    add_rect(slide, 0, Inches(0.08), SLIDE_W, Inches(1.05), PINE)
    if section_num:
        add_textbox(slide, Inches(0.55), Inches(0.22), Inches(1.2), Inches(0.35),
                    section_num, 11, True, OCHRE)
    add_textbox(slide, Inches(0.55), Inches(0.48), Inches(12), Inches(0.55),
                title, 26, True, WHITE)
    if subtitle:
        add_textbox(slide, Inches(0.55), Inches(0.95), Inches(12), Inches(0.35),
                    subtitle, 11, False, PAPER_DARK)
    add_textbox(slide, Inches(0.55), Inches(7.05), Inches(6), Inches(0.3),
                "頤安本草 · 網站重建整合方案", 9, False, MUTED)
    add_textbox(slide, Inches(10.5), Inches(7.05), Inches(2.3), Inches(0.3),
                f"{today.year}-06", 9, False, MUTED, PP_ALIGN.RIGHT)


def stat_card(slide, left, top, number, label, w=Inches(2.55), h=Inches(1.45)):
    add_rect(slide, left, top, w, h, WHITE)
    add_rect(slide, left, top, Inches(0.06), h, CINNABAR)
    add_textbox(slide, left + Inches(0.18), top + Inches(0.12), w - Inches(0.28), Inches(0.65),
                number, 32, True, PINE)
    add_textbox(slide, left + Inches(0.18), top + Inches(0.78), w - Inches(0.28), Inches(0.5),
                label, 10, False, MUTED)


def comparison_col(slide, left, top, width, header, header_color, items, bg=WHITE):
    h = Inches(0.42)
    add_rect(slide, left, top, width, h, header_color)
    add_textbox(slide, left + Inches(0.12), top + Inches(0.06), width - Inches(0.24), h,
                header, 12, True, WHITE)
    body_h = Inches(3.6)
    add_rect(slide, left, top + h, width, body_h, bg)
    add_bullets(slide, left + Inches(0.12), top + h + Inches(0.12),
                width - Inches(0.24), body_h - Inches(0.15), items, 10.5, INK)


def data_table(slide, y_start, headers, rows, col_widths=None):
    xs = [Inches(0.55)]
    if col_widths is None:
        col_widths = [Inches(2.0), Inches(3.2), Inches(3.2), Inches(3.2)]
    for w in col_widths[:-1]:
        xs.append(xs[-1] + w)
    row_h = Inches(0.62)
    # header
    y = y_start
    for i, h in enumerate(headers):
        add_rect(slide, xs[i], y, col_widths[i], Inches(0.38), PINE)
        add_textbox(slide, xs[i] + Inches(0.08), y + Inches(0.05), col_widths[i] - Inches(0.12),
                    Inches(0.3), h, 10, True, WHITE)
    y += Inches(0.38)
    for ri, row in enumerate(rows):
        bg = WHITE if ri % 2 == 0 else PAPER
        for ci, cell in enumerate(row):
            add_rect(slide, xs[ci], y, col_widths[ci], row_h, bg)
            add_textbox(slide, xs[ci] + Inches(0.08), y + Inches(0.06),
                        col_widths[ci] - Inches(0.14), row_h - Inches(0.1),
                        cell, 9.5, ci == 0, PINE if ci == 0 else INK)
        y += row_h


prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]

# ── 1 Title ──
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, PINE)
add_rect(s, 0, 0, Inches(0.12), SLIDE_H, CINNABAR)
add_textbox(s, Inches(0.9), Inches(1.85), Inches(8), Inches(0.45), "OAKVILLE WELLNESS", 13, False, OCHRE)
add_textbox(s, Inches(0.9), Inches(2.35), Inches(10), Inches(1.0),
            "頤安本草 · 官網重建整合方案", 38, True, WHITE)
add_textbox(s, Inches(0.9), Inches(3.45), Inches(10), Inches(0.55),
            "對標分析 · 三方比較 · 新站現況 · 廣告策略", 20, False, PAPER)
add_textbox(s, Inches(0.9), Inches(4.25), Inches(10), Inches(0.4),
            f"{today.year} 年 6 月 ｜ oakvilles.com ｜ 已部署 Vercel", 12, False, PAPER_DARK)

# ── 2 背景 ──
s = prs.slides.add_slide(blank)
slide_header(s, "01", "背景與一句話結論")
add_textbox(s, Inches(0.55), Inches(1.32), Inches(12.2), Inches(0.85),
            "舊站為單頁結構，無法承接「症狀搜尋」與「醫師名字搜尋」。新站借鑑養康的內容厚度與 WhatsApp 轉化，"
            "但保持中環高端、合規克制定位。",
            13, False, INK)
stat_card(s, Inches(0.55), Inches(2.35), "62", "新站頁面（中+英）")
stat_card(s, Inches(3.25), Inches(2.35), "7", "症狀 SEO 頁")
stat_card(s, Inches(5.95), Inches(2.35), "6", "dataLayer 事件")
stat_card(s, Inches(8.65), Inches(2.35), "✓", "Vercel 已上線")
add_rect(s, Inches(0.55), Inches(4.15), Inches(12.2), Inches(2.55), WHITE)
add_textbox(s, Inches(0.75), Inches(4.28), Inches(11.8), Inches(0.32), "簡報結構", 13, True, PINE)
add_bullets(s, Inches(0.75), Inches(4.65), Inches(11.5), Inches(1.9), [
    "三方對照（養康 / 舊站 / 新站）— 一表看清差距",
    "對標借鑑與品牌區隔 — 學架構、不學促銷",
    "新站轉化 + 量測 + Meta／Google 廣告落地路線",
], 11.5, INK)

# ── 3 三方對照（核心，取代舊「舊站不足」+「新舊改善」兩頁）──
s = prs.slides.add_slide(blank)
slide_header(s, "02", "三方多維度對照", "養康（標竿）｜舊站｜新站（2026.06 現況）")
data_table(s, Inches(1.42), ["維度", "養康 honscmc.com", "舊站", "新站（已上線）"], [
    ("定位", "尖沙咀 · 皮膚/美容 · 促銷", "個人醫師名 · 訊息散", "頤安本草 · 中環高端 · 合規"),
    ("頁面", "WP 多頁 · 四層 nav", "實質 1 頁", "62 頁（含 /en/）"),
    ("SEO", "每症狀一頁 + 地區詞", "無 sitemap/Schema", "7 症狀頁 + central-hk + OG"),
    ("轉化", "Joinchat + 首診優惠", "僅 WA 浮窗", "2 步 funnel + 計算器 + sticky"),
    ("信任", "13年/8000診量/FAQ", "缺評價/環境/FAQ", "Google 5.0 + gallery + 明碼價"),
    ("量度", "—", "無 GA4/CTA", "6 事件就緒 · GA4 待填"),
], [Inches(1.55), Inches(3.35), Inches(3.35), Inches(3.95)])

# ── 4 養康獲客（合併原設計+行銷兩頁）──
s = prs.slides.add_slide(blank)
slide_header(s, "03", "對標：養康獲客策略", "官網承接 + 社交種草 + 優惠降門檻")
comparison_col(s, Inches(0.55), Inches(1.42), Inches(5.9), "網站特徵", BENCH, [
    "WordPress 多頁 · 四層下拉",
    "綠色臨床 + 診所實景 Hero",
    "首屏：地區 + 專科 + WA 24h",
    "FAQ 10+ · Blog「中醫冷知識」",
    "Joinchat · IG 用家留言嵌入",
])
comparison_col(s, Inches(6.85), Inches(1.42), Inches(5.9), "獲客手段", CINNABAR, [
    "SEO 長尾：症狀 + 地區著陸頁",
    "WhatsApp 優先 · 全站立即預約",
    "首診優惠（皮膚免診金）",
    "IG KOL 素人化體驗分享",
    "PIXNET 體驗文 · 付費廣告著陸",
])
add_rect(s, Inches(0.55), Inches(5.85), Inches(12.2), Inches(0.95), PINE)
add_textbox(s, Inches(0.75), Inches(6.0), Inches(11.8), Inches(0.7),
            "頤安借鑑其「內容厚度 + WA 轉化」，但不照搬促銷語、免診金、論壇軟文。",
            11, False, PAPER)

# ── 5 借鑑 vs 不照搬 ──
s = prs.slides.add_slide(blank)
slide_header(s, "04", "借鑑 vs 不照搬")
comparison_col(s, Inches(0.55), Inches(1.42), Inches(5.9), "✓  借鑑", PINE, [
    "以「症狀」作 SEO 入口",
    "FAQ 結構化（地點、資歷、收費）",
    "WhatsApp 全站一致預約",
    "Blog 持續內容 + 內部連結",
    "首屏展示地點、時間、專科",
])
comparison_col(s, Inches(6.85), Inches(1.42), Inches(5.9), "✗  不照搬", CINNABAR, [
    "美容塑形促銷（豐胸、減肥）",
    "首診免診金（高端定位不符）",
    "efficacy 進取表述 · 論壇軟文",
    "資訊過密、綠色促銷臨床感",
    "與「循本溯源」定位衝突",
])

# ── 6 新站架構 ──
s = prs.slides.add_slide(blank)
slide_header(s, "05", "新站資訊架構", "Eleventy 靜態站 · dongfang.css · 雙語 /en/")
modules = [
    ("首頁", "Hero · 四專科 · gallery · 評價 · 計算器 · 預約"),
    ("專科/療法", "痛症·皮膚·婦科·內科 + 針灸·中藥·艾灸·拔罐"),
    ("症狀 SEO", "濕疹·暗瘡·失眠·備孕·頸痛·坐骨等 7 頁"),
    ("內容/信任", "Blog 4 篇 · FAQ · 醫師 · 流程收費 · Google 評價"),
    ("本地", "central-hk · GEO meta · llms.txt · hreflang"),
    ("部署", "GitHub → Vercel 自動 build · oakvilles.com"),
]
for i, (title, desc) in enumerate(modules):
    col = i % 3
    row = i // 3
    x = Inches(0.55) + col * Inches(4.15)
    y = Inches(1.48) + row * Inches(2.35)
    add_rect(s, x, y, Inches(3.85), Inches(2.05), WHITE)
    add_rect(s, x, y, Inches(3.85), Inches(0.4), PINE if row == 0 else PINE_LIGHT)
    add_textbox(s, x + Inches(0.12), y + Inches(0.05), Inches(3.5), Inches(0.32), title, 13, True, WHITE)
    add_textbox(s, x + Inches(0.12), y + Inches(0.52), Inches(3.5), Inches(1.4), desc, 11, False, INK)

# ── 7 品牌定位 ──
s = prs.slides.add_slide(blank)
slide_header(s, "06", "品牌定位差異")
pairs = [
    ("視覺", "綠色臨床、促銷感", "paper/pine/cinnabar 禪意"),
    ("受眾", "尖沙咀、美容減肥", "中環上班族、高端全科"),
    ("語氣", "促銷、療效導向", "循本溯源、合規克制"),
    ("品牌", "養康中醫館", "頤安本草 · 伍厚臻中醫師"),
]
y = Inches(1.58)
add_textbox(s, Inches(0.55), y, Inches(2.5), Inches(0.3), "維度", 11, True, PINE)
add_textbox(s, Inches(3.5), y, Inches(4.0), Inches(0.3), "養康", 11, True, MUTED)
add_textbox(s, Inches(7.8), y, Inches(4.5), Inches(0.3), "頤安新站", 11, True, CINNABAR)
y += Inches(0.4)
for i, (dim, old, new) in enumerate(pairs):
    bg = PAPER if i % 2 == 0 else WHITE
    add_rect(s, Inches(0.55), y, Inches(12.2), Inches(0.88), bg)
    add_textbox(s, Inches(0.7), y + Inches(0.18), Inches(2.5), Inches(0.5), dim, 12, True, PINE)
    add_textbox(s, Inches(3.5), y + Inches(0.18), Inches(3.8), Inches(0.5), old, 11, False, INK)
    add_textbox(s, Inches(7.8), y + Inches(0.18), Inches(4.5), Inches(0.5), new, 11, True, PINE)
    y += Inches(0.88)
add_rect(s, Inches(0.55), Inches(5.55), Inches(12.2), Inches(1.15), PINE)
add_textbox(s, Inches(0.75), Inches(5.72), Inches(11.8), Inches(0.85),
            "養康：流量 + 轉化優先 ｜ 頤安：信任 + 合規 + 高端體驗優先。架構可借鑑，調性刻意區隔。",
            12, False, PAPER)

# ── 8 轉化 + 量測 ──
s = prs.slides.add_slide(blank)
slide_header(s, "07", "轉化動線與量測", "WhatsApp 6734 9532 為統一 KPI")
add_rect(s, Inches(0.55), Inches(1.42), Inches(5.9), Inches(5.15), WHITE)
add_textbox(s, Inches(0.75), Inches(1.55), Inches(5.5), Inches(0.32), "轉化漏斗", 14, True, PINE)
add_bullets(s, Inches(0.75), Inches(1.95), Inches(5.5), Inches(4.4), [
    "多入口 CTA（hero / nav / sticky / 浮動 WA）",
    "每個按鈕帶 data-cta-id 可追蹤",
    "2 步預約：科別+日期 → 聯絡資料",
    "自動組 WhatsApp 訊息送出",
    "診金計算器 →「以此方案預約」預填",
    "手機 sticky bar + WA 浮動避障",
], 11.5, INK)
add_rect(s, Inches(6.85), Inches(1.42), Inches(5.9), Inches(5.15), WHITE)
add_textbox(s, Inches(7.05), Inches(1.55), Inches(5.5), Inches(0.32), "dataLayer 事件（GA4 就緒）", 14, True, CINNABAR)
add_bullets(s, Inches(7.05), Inches(1.95), Inches(5.5), Inches(4.4), [
    "whatsapp_click — 主要轉化 KPI",
    "booking_submit — 高意向次要轉化",
    "cta_click — 分辨入口（hero/sticky/wa-float）",
    "funnel_step — 預約第 1/2 步",
    "calculator_update / calculator_to_booking",
    "待辦：site-config.js 填入 GA4 ID + GTM",
], 11.5, INK)

# ── 9 行銷與廣告 ──
s = prs.slides.add_slide(blank)
slide_header(s, "08", "Meta / Google 廣告策略", "著陸頁對應 + 受眾 + 待接追蹤")
rows = [
    ("Meta 廣告", "症狀頁著陸 · Pixel 追蹤 WA", "濕疹/暗瘡/備孕 · /en/ 外籍"),
    ("Google Search", "品牌+中環+症狀長尾", "GA4 轉化匯入 · Maximize Conversions"),
    ("Google 本地", "GBP 優化 · 評價邀請", "與 Search 互補本地 pack"),
    ("IG 有機", "@oakville.wellness 互導", "診所日常 · 非硬銷（可選 KOL）"),
    ("內容 SEO", "FAQ + 症狀頁擴充", "Blog 每月 1–2 篇"),
    ("不建議", "免診金 · 論壇軟文", "豐胸減肥主 SEO 入口"),
]
y = Inches(1.48)
for i, (a, b, c) in enumerate(rows):
    bg = PAPER if i % 2 == 0 else WHITE
    add_rect(s, Inches(0.55), y, Inches(12.2), Inches(0.68), bg)
    add_textbox(s, Inches(0.7), y + Inches(0.1), Inches(2.4), Inches(0.48), a, 10.5, True, PINE)
    add_textbox(s, Inches(3.2), y + Inches(0.1), Inches(4.5), Inches(0.48), b, 10, False, INK)
    add_textbox(s, Inches(7.9), y + Inches(0.1), Inches(4.6), Inches(0.48), c, 10, False, MUTED)
    y += Inches(0.68)
add_rect(s, Inches(0.55), Inches(5.75), Inches(12.2), Inches(1.05), PINE_LIGHT)
add_textbox(s, Inches(0.75), Inches(5.9), Inches(11.8), Inches(0.75),
            "廣告開跑前 P0：GA4 ID · Search Console · GBP · GTM（GA4 + Meta Pixel）",
            11, False, PAPER)

# ── 10 交付路線圖 ──
s = prs.slides.add_slide(blank)
slide_header(s, "09", "交付路線圖（更新現況）")
phases = [
    ("第一階段", "62 頁雙語 · SEO · funnel · GA4 就緒", "✓ 已完成", PINE),
    ("第二階段", "50 張品牌視覺 · 全站接入", "規劃中 · 部分接入", OCHRE),
    ("第三階段", "Production oakvilles.com", "✓ 已部署 Vercel", PINE),
    ("持續", "月度行銷 HKD 10,000", "可選合約", CINNABAR),
]
for i, (phase, desc, status, color) in enumerate(phases):
    x = Inches(0.55) + i * Inches(3.15)
    add_rect(s, x, Inches(2.15), Inches(2.85), Inches(3.35), WHITE)
    add_rect(s, x, Inches(2.15), Inches(2.85), Inches(0.5), color)
    add_textbox(s, x + Inches(0.12), Inches(2.22), Inches(2.6), Inches(0.35), phase, 12, True, WHITE)
    add_textbox(s, x + Inches(0.12), Inches(2.85), Inches(2.6), Inches(1.7), desc, 10.5, False, INK)
    add_textbox(s, x + Inches(0.12), Inches(4.75), Inches(2.6), Inches(0.38), status, 11, True, color)
    if i < 3:
        arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + Inches(2.9), Inches(3.5), Inches(0.32), Inches(0.22))
        set_fill(arr, OCHRE)
        arr.line.fill.background()

# ── 11 下一步 ──
s = prs.slides.add_slide(blank)
slide_header(s, "10", "建議下一步")
cols = [
    ("P0 · 本週", PINE, [
        "填入 SITE_GA4_ID 並部署",
        "Search Console 提交 sitemap",
        "Google Business Profile 優化",
    ]),
    ("P1 · 廣告前", OCHRE, [
        "GTM：GA4 + Meta Pixel",
        "GA4 轉化 → Google Ads",
        "Meta 自訂轉換 whatsapp_click",
    ]),
    ("P2 · 持續", CINNABAR, [
        "症狀 campaign + 著陸 A/B",
        "Blog 每月 1–2 篇",
        "品牌視覺 50 張批次接入",
    ]),
]
for i, (title, color, items) in enumerate(cols):
    x = Inches(0.55) + i * Inches(4.15)
    add_rect(s, x, Inches(1.45), Inches(3.85), Inches(5.1), WHITE)
    add_rect(s, x, Inches(1.45), Inches(3.85), Inches(0.42), color)
    add_textbox(s, x + Inches(0.12), Inches(1.52), Inches(3.6), Inches(0.32), title, 13, True, WHITE)
    add_bullets(s, x + Inches(0.12), Inches(2.05), Inches(3.6), Inches(4.2), items, 11.5, INK)

# ── 12 結論 ──
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, PINE)
add_rect(s, 0, 0, Inches(0.12), SLIDE_H, CINNABAR)
add_textbox(s, Inches(0.9), Inches(1.45), Inches(11), Inches(0.4), "結論", 13, False, OCHRE)
add_textbox(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(2.2),
            "養康的優勢在於「內容厚度 + WhatsApp 轉化 + 社交種草」。"
            "舊站無法承接搜尋流量。新站已在架構、SEO、轉化量度及品牌視覺上系統補足，"
            "並保持中環高端、合規、專業定位。\n\n"
            "共同 KPI：whatsapp_click ｜ oakvilles.com",
            19, False, WHITE)
add_rect(s, Inches(0.9), Inches(4.85), Inches(11.5), Inches(1.75), PINE_LIGHT)
add_textbox(s, Inches(1.1), Inches(5.0), Inches(11), Inches(0.32), "立即行動", 12, True, OCHRE)
add_bullets(s, Inches(1.1), Inches(5.38), Inches(11), Inches(1.1), [
    "啟用 GA4 + 廣告追蹤（GTM / Pixel）",
    "GBP + 本地 SEO 持續優化",
    "症狀內容 + Blog 每月更新",
], 12, PAPER)

# ── 13 網絡行銷方案 ──
s = prs.slides.add_slide(blank)
slide_header(s, "11", "月度網絡行銷方案", "IG/FB 各 4 post · 官網 2 文 · Meta + Google 廣告")
stat_card(s, Inches(0.55), Inches(1.38), "4", "IG post／月", Inches(2.85), Inches(1.25))
stat_card(s, Inches(3.55), Inches(1.38), "4", "FB post／月", Inches(2.85), Inches(1.25))
stat_card(s, Inches(6.55), Inches(1.38), "2", "Blog 文章／月", Inches(2.85), Inches(1.25))
stat_card(s, Inches(9.55), Inches(1.38), "6734 9532", "統一轉化 KPI", Inches(2.85), Inches(1.25))
add_rect(s, Inches(0.55), Inches(2.85), Inches(5.9), Inches(3.85), WHITE)
add_textbox(s, Inches(0.75), Inches(2.98), Inches(5.5), Inches(0.3), "每月社交四大支柱", 12, True, PINE)
add_bullets(s, Inches(0.75), Inches(3.35), Inches(5.5), Inches(3.2), [
    "W1 症狀科普 → /conditions/",
    "W2 診所信任 → /about/ · Google 5.0",
    "W3 季節養生 → 當月 Blog",
    "W4 預約 CTA → WhatsApp 2 步 funnel",
], 10.5, INK)
add_rect(s, Inches(6.85), Inches(2.85), Inches(5.9), Inches(3.85), WHITE)
add_textbox(s, Inches(7.05), Inches(2.98), Inches(5.5), Inches(0.3), "Meta + Google 廣告", 12, True, CINNABAR)
add_bullets(s, Inches(7.05), Inches(3.35), Inches(5.5), Inches(3.2), [
    "Meta：症狀頁著陸 · Pixel · 再營銷",
    "Google：品牌+中環+症狀 Search · PMax",
    "預算參考：Meta 5–8K + Google 5–8K HKD/月",
    "P0：GA4 · GTM · Meta Pixel · GBP",
], 10.5, INK)
add_rect(s, Inches(0.55), Inches(6.85), Inches(12.2), Inches(0.45), PINE_LIGHT)
add_textbox(s, Inches(0.75), Inches(6.92), Inches(11.8), Inches(0.32),
            "不合規手段不採用：免診金 · 論壇軟文 · 豐胸減肥主 SEO ｜ 詳見 Oakville-Digital-Marketing-Plan.docx",
            9.5, False, PAPER)

prs.save(OUTPUT)
shutil.copy2(OUTPUT, OUTPUT_DESKTOP)
os.makedirs(os.path.dirname(OUTPUT_PRESENTATIONS), exist_ok=True)
shutil.copy2(OUTPUT, OUTPUT_PRESENTATIONS)
print(f"Saved: {OUTPUT}")
print(f"Copied: {OUTPUT_DESKTOP}")
print(f"Copied: {OUTPUT_PRESENTATIONS}")
print(f"Slides: {len(prs.slides)}")
